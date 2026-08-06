"""
build_warehouse_pptx.py — 倉管 Agent v2 專案簡報（給老闆看，圖文並茂）
風格：白底 + Microsoft JhengHei（中文）+ Arial（英文/數字/code）+ mint teal
字級一律 ≥ 12 pt（頁碼/註腳除外）
三段式：產品價值 → 系統架構 → 測試方法
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION

# ─── Theme ──────────────────────────────────────────────────
TEAL    = RGBColor(0x00, 0xC4, 0x9A)   # 主色 mint teal
TEALDK  = RGBColor(0x00, 0x8E, 0x76)   # 深一階
DARK    = RGBColor(0x1A, 0x1A, 0x2E)   # 深藍黑（標題/深底）
NAVY    = RGBColor(0x22, 0x2A, 0x45)
GREY44  = RGBColor(0x44, 0x44, 0x44)
GREY55  = RGBColor(0x55, 0x55, 0x55)
GREY77  = RGBColor(0x77, 0x77, 0x77)
GREYBB  = RGBColor(0xBB, 0xBB, 0xBB)
GREYE6  = RGBColor(0xE6, 0xE8, 0xEC)
LIGHT   = RGBColor(0xF4, 0xF6, 0xF8)   # 卡片淺底
TEALBG  = RGBColor(0xE4, 0xF7, 0xF2)   # 主色淺底
AMBER   = RGBColor(0xF0, 0xA0, 0x30)
CORAL   = RGBColor(0xE8, 0x5D, 0x5D)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

FONT_ZH = "Microsoft JhengHei"
FONT_EN = "Arial"

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
MX = Inches(0.73)   # 左右邊界

OUT = Path(__file__).parent / "倉管Agent_專案簡報.pptx"

# ─── 語音模型：單一開關 ──────────────────────────────────────
# 2026-08-06：q5_0 與 q8_0 二選一中。模型名稱／體積／延遲散在 13 處，
#   集中在這裡，定案後只要改 ASR_PICK 一行，全簡報自動一致。
#   數字全部來自 2026-08-06 同一套實測（英文版·機二·網頁真實路徑·只計 ASR 段）。
ASR_PICK = "q8_0"          # 2026-08-06 定案（端到端 78.0% 勝 q5 的 77.3%，且快 40%）
_ASR_SPEC = {
    "q5_0": {"file": "ggml-small-q5_0.bin", "name": "small-q5_0",
             "size": "167 MB", "bits": "壓最小（5-bit）", "lat": "4.2s",
             "e2e": "76.7%"},
    "q8_0": {"file": "ggml-small-q8_0.bin", "name": "small-q8_0",
             "size": "252 MB", "bits": "壓中等（8-bit）", "lat": "2.5s",
             "e2e": "78.0%"},
}
ASR = _ASR_SPEC[ASR_PICK]

# ─── 守衛庫句數：單一來源 ────────────────────────────────────
# 2026-08-06 補課：ZH 1122→1149、EN 892→938。
#   （守衛庫自 7/20（ZH）／7/25（EN）停止成長，期間十餘輪修復未納入回歸，
#     今日補齊；數字散在 10 處，集中管理避免下次又對不上。）
GUARD_ZH = 1149
GUARD_EN = 938


# ─── Helpers ─────────────────────────────────────────────────
def set_run(run, *, text=None, font=FONT_ZH, size=14, bold=False, color=DARK):
    if text is not None:
        run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(slide, x, y, w, h, text, *, font=FONT_ZH, size=14, bold=False,
             color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.15):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Emu(0))
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else [str(text)]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        set_run(run, text=line, font=font, size=size, bold=bold, color=color)
    return box


def add_rich(slide, x, y, w, h, segments, *, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Emu(0))
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = [segments] if (segments and isinstance(segments[0], dict)) else segments
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for seg in line:
            run = p.add_run()
            set_run(run, text=seg.get("text", ""), font=seg.get("font", FONT_ZH),
                    size=seg.get("size", 14), bold=seg.get("bold", False),
                    color=seg.get("color", DARK))
    return box


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    if shadow:
        _soft_shadow(shp)
    return shp


def add_round(slide, x, y, w, h, fill=None, line=None, line_w=0.75, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    try:
        shp.adjustments[0] = 0.08
    except Exception:
        pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    if shadow:
        _soft_shadow(shp)
    return shp


def add_circle(slide, x, y, d, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    shp.shadow.inherit = False
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    return shp


def add_icon_circle(slide, x, y, d, glyph, *, circle=TEAL, gcolor=WHITE, gsize=18):
    """圓形色底 + 置中字元圖示（自動把 emoji 換成 ICON 表的中文字，跨平台一致）"""
    ch = ICON.get(glyph, glyph)
    add_circle(slide, x, y, d, circle)
    add_text(slide, x, y - Inches(0.02), d, d, ch, font=FONT_ZH, size=gsize,
             bold=True, color=gcolor, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# emoji → 乾淨中文字圖示（放進色圓，跨機器 100% 一致，比 emoji 專業）
ICON = {
    "🔍": "查", "⌨️": "鍵", "📱": "機", "💸": "$",
    "💬": "問", "⚡": "秒", "📶": "碼", "🍓": "Pi",
    "📦": "庫", "🔄": "動", "⚠️": "警", "⚖️": "比",
    "🔥": "熱", "⏰": "期", "🔗": "配", "🤖": "AI",
    "🗣️": "語", "✏️": "改", "🧭": "路", "🛡️": "盾", "⚙️": "行",
    "🎯": "準", "💰": "省", "🔌": "電", "�Ⓜ": "M",
    "🧠": "腦", "🔁": "環", "♻️": "重", "🔒": "鎖", "💡": "！",
    "✅": "✓",
}


def dot_icon(slide, x, y, emoji, *, d=0.42, circle=TEAL, gcolor=WHITE, gsize=13):
    """在 (x,y) 放一個小色圓 + ICON 字，取代裸 emoji。回傳圓的右緣 x 供接文字。"""
    ch = ICON.get(emoji, emoji)
    add_circle(slide, x, y, Inches(d), circle)
    add_text(slide, x, y - Inches(0.015), Inches(d), Inches(d), ch,
             font=FONT_ZH, size=gsize, bold=True, color=gcolor,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return x + Inches(d)


def add_arrow(slide, x, y, w, h, fill=TEAL):
    shp = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, x, y, w, h)
    shp.shadow.inherit = False
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp


def _soft_shadow(shp):
    """柔和陰影。effectLst 在 spPr 內必須排在 <a:ln> 之後（OOXML 子元素順序嚴格，
    位置錯 PowerPoint 會判檔案損壞、要求修復）。用 insert-after-ln 確保順序正確，
    不用 SubElement 盲目 append。"""
    from lxml import etree
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    spPr = shp._element.spPr
    # 已有 effectLst 就不重複加
    if spPr.find(f"{{{ns}}}effectLst") is not None:
        return
    effLst = etree.SubElement(spPr, f"{{{ns}}}effectLst")
    outer = etree.SubElement(effLst, f"{{{ns}}}outerShdw")
    outer.set("blurRad", "50800"); outer.set("dist", "25400")
    outer.set("dir", "5400000"); outer.set("rotWithShape", "0")
    clr = etree.SubElement(outer, f"{{{ns}}}srgbClr"); clr.set("val", "1A1A2E")
    alpha = etree.SubElement(clr, f"{{{ns}}}alpha"); alpha.set("val", "18000")
    # 移到正確位置：effectLst 必須在 ln 之後（若有 ln），否則保持在末尾
    ln = spPr.find(f"{{{ns}}}ln")
    if ln is not None:
        spPr.remove(effLst)
        ln.addnext(effLst)


def slide_blank():
    return prs.slides.add_slide(prs.slide_layouts[6])


def title_bar(slide, kicker, title):
    """章節頁首：小 kicker + 大標題（無底線，用留白區隔）"""
    add_text(slide, MX, Inches(0.34), Inches(11.8), Inches(0.30),
             kicker, font=FONT_EN, size=12.5, bold=True, color=TEAL)
    add_text(slide, MX, Inches(0.64), Inches(11.8), Inches(0.62),
             title, size=27, bold=True, color=DARK)


_PN = [0]
def pn(slide):
    _PN[0] += 1
    page_num(slide, _PN[0])


def page_num(slide, n):
    add_text(slide, Inches(12.55), Inches(7.18), Inches(0.72), Inches(0.26),
             f"{n:02d}", font=FONT_EN, size=10, color=GREYBB, align=PP_ALIGN.RIGHT)


def set_notes(slide, notes):
    nt = slide.notes_slide.notes_text_frame
    nt.text = notes
    for para in nt.paragraphs:
        for run in para.runs:
            run.font.name = FONT_ZH
            run.font.size = Pt(12)


def kpi_row(slide, top, items, *, box_h=1.18, num_size=30, on_dark=False):
    """items: (number, label)；淺卡片 + 大數字"""
    n = len(items)
    gap = Inches(0.22)
    total = SLIDE_W - 2 * MX
    bw = (total - gap * (n - 1)) / n
    for i, (num, lab) in enumerate(items):
        x = MX + (bw + gap) * i
        bg = NAVY if on_dark else LIGHT
        add_round(slide, x, top, bw, Inches(box_h), fill=bg, shadow=not on_dark)
        add_text(slide, x + Inches(0.18), top + Inches(0.16),
                 bw - Inches(0.36), Inches(0.60), num, font=FONT_EN,
                 size=num_size, bold=True, color=TEAL, align=PP_ALIGN.LEFT)
        add_text(slide, x + Inches(0.18), top + Inches(box_h) - Inches(0.42),
                 bw - Inches(0.36), Inches(0.34), lab, font=FONT_ZH,
                 size=12.5, color=(GREYBB if on_dark else GREY55),
                 align=PP_ALIGN.LEFT)


print("helpers ready")


# ═══════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ─── S1 封面（深底）─────────────────────────────────────────
s = slide_blank()
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=DARK)
add_rect(s, 0, 0, SLIDE_W, Inches(0.14), fill=TEAL)
add_text(s, MX, Inches(1.75), Inches(11.8), Inches(0.4),
         "邊緣 AI 應用展示 · 晶片團隊軟體成果", font=FONT_ZH, size=17, bold=True, color=TEAL)
add_text(s, MX, Inches(2.25), Inches(11.8), Inches(1.5),
         "倉管 Agent v2", size=54, bold=True, color=WHITE)
add_text(s, MX, Inches(3.35), Inches(11.8), Inches(0.7),
         "一句自然語言 → 秒級拿到正確答案，全程在一台 Raspberry Pi 上跑",
         size=18, color=GREYBB)
# 三個亮點膠囊
caps = [("270M", "邊緣級小模型"), ("100+ 輪", "品質收斂"), ("100%", "雙平台回歸")]
cw = Inches(2.75); cy = Inches(4.7); gap = Inches(0.3)
for i, (a, b) in enumerate(caps):
    x = MX + (cw + gap) * i
    add_round(s, x, cy, cw, Inches(1.15), fill=NAVY)
    add_text(s, x, cy + Inches(0.18), cw, Inches(0.55), a, font=FONT_EN,
             size=26, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    add_text(s, x, cy + Inches(0.72), cw, Inches(0.32), b, size=13,
             color=GREYBB, align=PP_ALIGN.CENTER)
add_text(s, MX, Inches(6.85), Inches(11.8), Inches(0.35),
         "FunctionGemma 270M 微調模型 × 多層校正架構 × RPi5 CPU 部署 → 下一步接自研晶片加速",
         font=FONT_EN, size=12, color=GREY77)
set_notes(s, "封面。定位：這是晶片團隊的軟體成果——一個能展示邊緣 AI 落地的真實應用。"
             "一句話說清楚產品：自然語言問倉管、秒級回答、跑在便宜的樹莓派上。"
             "三個數字：270M 是模型只有 2.7 億參數（業界主流 3B-8B 的零頭）、"
             "100+ 輪是品質收斂的迭代次數、100% 是回歸測試通過率。"
             "底線一句話點出硬體路線圖：現在純 CPU、下一步接自研晶片加速。")
pn(s)

# ─── S2 痛點 → 解法 ─────────────────────────────────────────
s = slide_blank()
title_bar(s, "THE PROBLEM", "倉管人員要的，不是又一套要學的系統")
# 左：傳統
add_round(s, MX, Inches(1.7), Inches(5.75), Inches(4.6), fill=LIGHT, shadow=True)
add_text(s, MX + Inches(0.35), Inches(1.95), Inches(5.05), Inches(0.5),
         "傳統倉管系統", size=18, bold=True, color=GREY44)
pains = [("🔍", "要記選單在哪、欄位怎麼填"),
         ("⌨️", "查個庫存點五六層才到"),
         ("📱", "現場拿手機打字，介面難用"),
         ("💸", "客製一次動輒數十萬")]
for i, (ic, tx) in enumerate(pains):
    y = Inches(2.65) + Inches(0.85) * i
    dot_icon(s, MX + Inches(0.35), y, ic, circle=GREYBB, gcolor=WHITE)
    add_text(s, MX + Inches(1.0), y + Inches(0.07), Inches(4.4), Inches(0.6),
             tx, size=14.5, color=GREY44, anchor=MSO_ANCHOR.MIDDLE)
# 右：本專案
rx = Inches(6.85)
add_round(s, rx, Inches(1.7), Inches(5.75), Inches(4.6), fill=TEALBG, shadow=True)
add_text(s, rx + Inches(0.35), Inches(1.95), Inches(5.05), Inches(0.5),
         "倉管 Agent v2", size=18, bold=True, color=TEALDK)
gains = [("💬", "直接用講的：「藍牙耳機還剩幾個」"),
         ("⚡", "一句話秒級回答，免選單"),
         ("📶", "掃 QR 手機就能問，展場離線也行"),
         ("🍓", "一台樹莓派跑得動，硬體成本極低")]
for i, (ic, tx) in enumerate(gains):
    y = Inches(2.65) + Inches(0.85) * i
    dot_icon(s, rx + Inches(0.35), y, ic, circle=TEAL, gcolor=WHITE)
    add_text(s, rx + Inches(1.0), y + Inches(0.07), Inches(4.4), Inches(0.6),
             tx, size=14.5, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
set_notes(s, "痛點對照。左邊是傳統倉管系統的四個痛：要學、層級深、現場難用、客製貴。"
             "右邊是本專案的解法：用講的、秒級答、手機掃碼、樹莓派就能跑。")
pn(s)

print("S1-S2 done")


# ─── S3 七大功能（2x4 grid + icon 圓）──────────────────────
s = slide_blank()
title_bar(s, "WHAT IT DOES", "訪客能問什麼：七大查詢 + 三類進階")
funcs = [
    ("📦", "庫存查詢", "商品 / 倉庫 / 類別"),
    ("🔄", "進出紀錄", "今天 / 本週 / 本月"),
    ("⚠️", "缺貨警示", "撐幾天 + 建議補多少"),
    ("⚖️", "倉庫比較", "任兩倉對比"),
    ("🔥", "熱銷排行", "期間 / 類別 TOP 10"),
    ("⏰", "到期預警", "N 天內即將過期"),
    ("🔗", "相關推薦", "購物籃搭售分析"),
    ("🤖", "Agent 進階", "根因分析 / 補貨 / 排程"),
]
cols, rows = 4, 2
gw = Inches(2.92); gh = Inches(2.0)
gapx = Inches(0.18); gapy = Inches(0.30)
x0 = MX; y0 = Inches(1.75)
for i, (ic, name, desc) in enumerate(funcs):
    r, c = divmod(i, cols)
    x = x0 + (gw + gapx) * c
    y = y0 + (gh + gapy) * r
    last = (i == 7)
    add_round(s, x, y, gw, gh, fill=(TEALBG if last else LIGHT), shadow=True)
    add_icon_circle(s, x + Inches(0.28), y + Inches(0.26), Inches(0.72), ic,
                    circle=(TEALDK if last else TEAL), gsize=20)
    add_text(s, x + Inches(0.28), y + Inches(1.12), gw - Inches(0.5), Inches(0.4),
             name, size=16, bold=True, color=DARK)
    add_text(s, x + Inches(0.28), y + Inches(1.52), gw - Inches(0.5), Inches(0.4),
             desc, size=12, color=GREY55)
set_notes(s, "七大查詢工具（唯讀）加一組 Agent 進階工具。前七個是訪客最常問的；"
             "第八格 Agent 進階包含根因分析（帳對不上時自動查原因）、自動補貨、定時排程。")
pn(s)

# ─── S4 一句話 → 一張卡（範例對照）──────────────────────────
s = slide_blank()
title_bar(s, "HOW IT FEELS", "訪客怎麼問，系統怎麼答")
examples = [
    ("藍牙耳機還剩幾個", "三倉共 227 件（北 141 / 中 42 / 南 44）", "📦"),
    ("哪些快沒貨了", "42 項低於安全庫存，最急：藍牙喇叭剩 15 件、12 天斷貨", "⚠️"),
    ("北倉進 50 個滑鼠", "跳出確認卡 → 按鈕確認才寫入庫存", "🔄"),
    ("電動牙刷的帳怎麼兜不起來", "自動查 PO → 找到短收 → 建議聯絡供應商", "🤖"),
]
y0 = Inches(1.8)
for i, (q, a, ic) in enumerate(examples):
    y = y0 + Inches(1.24) * i
    # 問句（左，深色膠囊）
    add_round(s, MX, y, Inches(4.75), Inches(1.02), fill=DARK, shadow=True)
    dot_icon(s, MX + Inches(0.28), y + Inches(0.3), "💬", circle=TEAL, gcolor=DARK)
    add_text(s, MX + Inches(0.88), y + Inches(0.12), Inches(3.7), Inches(0.78),
             q, size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    # 箭頭
    add_arrow(s, Inches(5.65), y + Inches(0.33), Inches(0.55), Inches(0.36), fill=TEAL)
    # 回答（右，淺卡）
    add_round(s, Inches(6.4), y, Inches(6.2), Inches(1.02), fill=LIGHT, shadow=True)
    dot_icon(s, Inches(6.65), y + Inches(0.3), ic, circle=TEALDK, gcolor=WHITE)
    add_text(s, Inches(7.25), y + Inches(0.12), Inches(5.1), Inches(0.82),
             a, size=13, color=GREY44, anchor=MSO_ANCHOR.MIDDLE)
set_notes(s, "四組真實對話。第三個示範「寫入類」操作有 HITL 確認卡——不會打字就直接改庫存，"
             "一定要按按鈕，避免誤操作。第四個是 Agent 根因分析：不只查數字，會推理為什麼對不上。")
pn(s)

print("S3-S4 done")


# ─── S4a Agent 三層自動化（證明是 Agent 不是查詢工具）★ ───────────
s = slide_blank()
title_bar(s, "AGENT · 三層自動化", "不是查詢介面，是會自己做事的 Agent")
add_text(s, MX, Inches(1.42), Inches(11.8), Inches(0.4),
         "感知 → 推理 → 行動，業界對 AI Agent 的三個標準能力，本系統三層都真的實作、都在背景跑。",
         size=13.5, color=GREY55)
auto_tiers = [
    ("感知 · PERCEIVE", "系統自己盯，不用人問", TEALDK,
     ["缺貨警示：背景每小時掃描，低於安全庫存主動推播",
      "到期預警：自動追蹤保存期限，N 天內主動提醒"]),
    ("推理 · REASON", "系統自己判斷該怎麼辦", TEAL,
     ["根因分析：帳兜不攏 → 自動掃 PO → 找短收 → 給建議（3 步推理）",
      "自動補貨：算出撐幾天、該補多少，產採購單草稿"]),
    ("行動 · ACT", "系統自己動手做", NAVY,
     ["定時排程：「每天 8 點跑盤點」背景每分鐘檢查、到點自動執行",
      "一句話改庫存：自然語言 → 抽參數 → 執行（10 種 HITL 確認把關）"]),
]
y0 = Inches(2.0); bh = Inches(1.42); bgap = Inches(0.18)
for i, (tier, sub, col, items) in enumerate(auto_tiers):
    y = y0 + (bh + bgap) * i
    add_round(s, MX, y, Inches(3.35), bh, fill=col)
    add_text(s, MX + Inches(0.28), y + Inches(0.3), Inches(2.85), Inches(0.4),
             tier, font=FONT_EN, size=14.5, bold=True, color=WHITE)
    add_text(s, MX + Inches(0.28), y + Inches(0.78), Inches(2.85), Inches(0.4),
             sub, size=12.5, color=WHITE)
    add_round(s, Inches(4.25), y, Inches(8.35), bh, fill=LIGHT, shadow=True)
    for j, it in enumerate(items):
        yy = y + Inches(0.28) + Inches(0.56) * j
        dot_icon(s, Inches(4.5), yy, "✅", d=0.34, circle=col, gcolor=WHITE, gsize=11)
        add_text(s, Inches(4.98), yy + Inches(0.01), Inches(7.4), Inches(0.44),
                 it, size=13, color=GREY44, anchor=MSO_ANCHOR.MIDDLE)
set_notes(s, "★Agent 定位頁。核心論點：這不是查詢工具，是會自己做事的 Agent。"
             "業界對 AI Agent 的標準定義＝感知→推理→行動三個能力，本系統三層都真的實作、"
             "背景在跑。感知：警示每小時掃、到期預警主動提醒。推理：RCA 三步根因分析（不只查"
             "現況，還找原因給建議）、自動補貨算量。行動：排程每分鐘檢查自動執行、一句話改庫存。"
             "分寸：HITL 確認卡不是缺點是優點——Agent 提案、人拍板，符合企業對 AI 落地的安全要求。")
pn(s)


# ─── S12c 動態倉庫模擬 Live 模式 ★（2026-08-03 新增）──────────────
s = slide_blank()
title_bar(s, "LIVE WAREHOUSE · 動態模擬", "倉庫自己會動：一天的進出貨濃縮成幾分鐘")
add_text(s, MX, Inches(1.40), Inches(11.8), Inches(0.62),
         "展場訪客盯著凍結快照看會覺得假。業界真實架構本來就是 perpetual inventory——庫存由條碼槍、"
         "RFID、電商訂單多來源即時更新，AI 助理是其上的對話層。模擬不是特效，是把真實架構演出來。",
         size=13, color=GREY55)
kpi_row(s, Inches(2.18), [
    ("200×", "時間加速（現場可調 1–400×）"),
    ("2.7s", "一輪背景進出貨寫入"),
    ("79%", "出庫占比＝seed 真值比例"),
    ("0", "對測試的干擾（守衛自動隔離）"),
], box_h=1.3, num_size=34)
_live_pts = [
    ("真實來源", "背景寫入掛 pda_scan / wms_sync / ecom_order 三種 actor——訪客查異動看到"
                 "混合來源，正好證明 Agent 看得到整個倉庫，不是只看得到自己"),
    ("護欄", "庫存在安全線 0.8–1.6 倍帶內波動、任何倉不低於 5 件——連跑三天不壞、"
             "reset 一鍵回乾淨基準（開機自動歸零，斷電/離線都免時鐘）"),
    ("誠實的代價", "模擬灌大資料曾把 CPU 燒滿（py-spy 抓到 /anomalies 每次輪詢全掃數十萬筆）"
                   "——以 30 秒快取＋出貨日索引三層修復，CPU 200%→31%、查詢照常秒回"),
    ("數據紀律", "「昨天/報表/撐天」的統計永遠只取乾淨歷史——模擬寫入被分析層隔離，"
                 "看得到即時跳動、算得出正確數字，兩者不打架"),
]
for i, (k, v) in enumerate(_live_pts):
    y = Inches(3.78) + Inches(0.80) * i
    dot_icon(s, MX, y + Inches(0.05), "●", d=0.34, circle=TEAL, gsize=10)
    add_text(s, MX + Inches(0.55), y, Inches(1.9), Inches(0.7), k, size=13.5, bold=True,
             color=TEALDK)
    add_text(s, MX + Inches(2.5), y, Inches(9.3), Inches(0.74), v, size=12.5, color=GREY44)
set_notes(s, "★動態模擬頁（2026-08-03 上線）。訪客第一眼就看到三倉數字在跳、告警橫幅自己出現"
             "——回應「倉庫不動很假」的展場回饋。查證過業界：現代倉儲就是 perpetual inventory，"
             "多數庫存變動沒有人對系統下指令（條碼槍/RFID/電商自動寫入），對話式 AI 是其上的"
             "查詢層——所以這不是演假資料，是把真實架構濃縮演出。工程重點四件事：①actor 不偽裝"
             "訪客操作，查異動看得到混合來源；②護欄讓它跑三天不壞；③效能課金誠實講：模擬灌出"
             "數十萬筆後，前端每 8 秒輪詢的異常掃描把兩核燒滿（py-spy 實錘），用 TTL 快取＋"
             "出貨日索引修復——這段也是很好的工程故事；④統計紀律：分析層一律排除模擬寫入，"
             "『昨天』永遠是乾淨歷史，數字對得起對帳。")
pn(s)


# ─── S4b OOV 招牌能力（訪客怎麼亂打都聽得懂）★重點 ─────────────
s = slide_blank()
title_bar(s, "SIGNATURE · OOV 容錯", "訪客怎麼亂打，系統照樣聽懂")
add_text(s, MX, Inches(1.42), Inches(11.8), Inches(0.4),
         "展場訪客不會照規矩打字——錯字、注音殘字、講一半、中英夾雜、講俗稱。"
         "270M 小模型 + 多層容錯，照樣抓對商品。", size=13.5, color=GREY55)
oov_cards = [
    ("錯字 / 同音", "「藍芽耳機」「汽泡水」\n「悶燒灌」「瑜加墊」", "→ 自動修正抓對商品"),
    ("講一半 / 不完整", "「那個充電的」\n「洗衣的那個精」", "→ 補全成完整商品名"),
    ("注音殘字", "「安全ㄎ存多少」\n「那ㄍ快到期嗎」", "→ 注音還原成國字"),
    ("中英夾雜", "「藍牙 earphone 庫存」\n「coffee 豆剩多少」", "→ 中英混打照抓"),
    ("語序顛倒 / 口語", "「庫存藍牙耳機」\n「剩多少氣泡水啊」", "→ 不管詞序都懂"),
    ("俗稱 / 別名", "「充手機的寶」\n「做瑜伽的墊子」", "→ 對應到正式品名"),
]
cols = 3
gw = Inches(3.87); gh = Inches(2.15); gapx = Inches(0.13); gapy = Inches(0.22)
x0 = MX; y0 = Inches(2.05)
for i, (cat, ex, res) in enumerate(oov_cards):
    r, c = divmod(i, cols)
    x = x0 + (gw + gapx) * c
    y = y0 + (gh + gapy) * r
    add_round(s, x, y, gw, gh, fill=LIGHT, shadow=True)
    add_text(s, x + Inches(0.25), y + Inches(0.2), gw - Inches(0.5), Inches(0.4),
             cat, size=14.5, bold=True, color=TEALDK)
    add_text(s, x + Inches(0.25), y + Inches(0.68), gw - Inches(0.5), Inches(0.9),
             ex, size=13.5, color=DARK, line_spacing=1.25)
    add_text(s, x + Inches(0.25), y + Inches(1.68), gw - Inches(0.5), Inches(0.38),
             res, size=12.5, color=TEAL, anchor=MSO_ANCHOR.MIDDLE)
set_notes(s, "★招牌能力頁。這是軟體最難、最值得秀的地方：小模型還能扛住這些亂打。"
             "六類真實展場輸入——錯字、講一半、注音殘字、中英夾雜、語序顛倒、俗稱別名。"
             "全部抓對商品。這些例子都是實測語料裡的真句子，不是虛構。")
pn(s)


# ─── S4c OOV 怎麼做到的（多層容錯架構）─────────────────────────
s = slide_blank()
title_bar(s, "HOW · 多層容錯", "小模型抓不準？多層防線把它接住")
add_text(s, MX, Inches(1.42), Inches(11.8), Inches(0.4),
         "270M 抽出的 keyword 常常髒、殘、錯。四層 fallback 逐層修，抓不到才反問——"
         "絕不亂猜、絕不幻覺出錯的商品。", size=13.5, color=GREY55)
oov_layers = [
    ("1", "雜詞剝除", "剝掉倉庫名 / 量詞 / 語氣詞 / 功能詞尾巴", "北倉的藍芽耳機還有幾個  →  藍芽耳機", NAVY),
    ("2", "精確 / 錯字修復", "同音錯字表 + 注音還原，換成正確品名", "藍芽耳機  →  無線藍牙耳機", TEALDK),
    ("3", "模糊比對 fuzzy", "剝規格 + 雙向滑窗 + 字元重疊，容忍變體", "充手機的寶  →  行動電源 10000mAh", TEAL),
    ("4", "抓不到就反問", "列疑似商品請選，附庫存概況，絕不亂猜", "帽子  →  你是指遮陽帽還是毛帽？", AMBER),
]
y0 = Inches(2.05); bh = Inches(1.02); bgap = Inches(0.16)
for i, (num, name, desc, ex, col) in enumerate(oov_layers):
    y = y0 + (bh + bgap) * i
    add_round(s, MX, y, Inches(0.82), bh, fill=col)
    add_text(s, MX, y, Inches(0.82), bh, num, font=FONT_EN, size=26, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_round(s, Inches(1.7), y, Inches(10.9), bh, fill=LIGHT, shadow=True)
    add_text(s, Inches(1.95), y + Inches(0.13), Inches(2.9), Inches(0.4),
             name, size=15.5, bold=True, color=col)
    add_text(s, Inches(1.95), y + Inches(0.55), Inches(2.9), Inches(0.4),
             desc, size=11, color=GREY55)
    add_text(s, Inches(5.05), y + Inches(0.13), Inches(7.35), Inches(0.76),
             ex, font=FONT_ZH, size=13, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
set_notes(s, "OOV 容錯的技術原理：四層 fallback。①雜詞剝除②精確/錯字修復（同音錯字表+"
             "注音還原）③模糊比對 fuzzy（剝規格、雙向滑窗、字元重疊，容忍變體詞）"
             "④抓不到就反問（列疑似商品請選，絕不亂猜幻覺）。核心哲學：寧可反問，不可答錯。")
pn(s)


# ─── S5 架構分層（縱向流程圖）──────────────────────────────
s = slide_blank()
title_bar(s, "ARCHITECTURE", "小模型只做路由，業務邏輯交給伺服器")
add_text(s, MX, Inches(1.42), Inches(11.8), Inches(0.4),
         "設計理念：270M 小模型負責「意圖分類」，Server 端負責精確計算與業務規則——"
         "讓 LLM 只做它擅長的事", size=13.5, color=GREY55)
layers = [
    ("使用者輸入", "口語自然語言（含錯字 / 注音 / 英文俗稱）", GREY44, "🗣️"),
    ("Query Rewriting", "50+ 條規則（6 大類）把口語 → 標準句型", NAVY, "✏️"),
    ("intent_clf 主路由", "FastText 14 類意圖分類，先決定用哪個功能（98.9% 準，下頁詳介）", TEALDK, "🧭"),
    ("FunctionGemma 270M", "小模型只抽參數：商品 / 倉庫 / 時間", TEAL, "🤖"),
    ("dispatch + 校正層", "C0–C18 規則層：最後防線，擋幻覺、補接地", NAVY, "🛡️"),
    ("業務工具執行", "查詢 / 庫存異動 / 根因分析 / 排程", DARK, "⚙️"),
]
y0 = Inches(2.0); bh = Inches(0.72); bgap = Inches(0.14)
for i, (name, desc, col, ic) in enumerate(layers):
    y = y0 + (bh + bgap) * i
    add_round(s, MX, y, Inches(0.72), bh, fill=col)
    add_text(s, MX, y, Inches(0.72), bh, ICON.get(ic, ic), size=18, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_round(s, Inches(1.6), y, Inches(11.0), bh, fill=LIGHT, shadow=True)
    add_text(s, Inches(1.9), y + Inches(0.09), Inches(4.0), Inches(0.55),
             name, font=FONT_EN, size=15, bold=True, color=col,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(5.7), y + Inches(0.09), Inches(6.7), Inches(0.55),
             desc, size=13, color=GREY44, anchor=MSO_ANCHOR.MIDDLE)
    if i < len(layers) - 1:
        add_text(s, Inches(0.86), y + bh - Inches(0.04), Inches(0.45), Inches(0.2),
                 "▼", font=FONT_EN, size=10, color=GREYBB, align=PP_ALIGN.CENTER)
set_notes(s, "六層架構，由上而下。核心洞見：270M 小模型輸出不穩定是先天限制，"
             "所以不讓它做決策，只做意圖分類與抽參數；真正的業務邏輯與精確計算放在"
             "Server 端的規則層（C0-C18）。這也是業界邊緣 Agent 的正解——小模型當路由器，不當決策者。")
pn(s)

# ─── S5b 小主角 FastText intent_clf ────────────────────────
s = slide_blank()
title_bar(s, "THE QUIET HERO", "小主角 FastText：路由交給穩的，生成交給會的")
kpi_row(s, Inches(1.52), [
    ("14 類", "自定義意圖標籤"),
    ("4.2MB", "intent_clf 量化模型（原 0.5GB）"),
    ("98.9%", "路由準確率"),
    ("毫秒級", "CPU 路由延遲"),
], num_size=27)

# 左卡：為什麼主路由是它
_ly = Inches(2.98)
add_round(s, MX, _ly, Inches(6.55), Inches(3.72), fill=LIGHT, shadow=True)
add_text(s, MX + Inches(0.35), _ly + Inches(0.22), Inches(5.8), Inches(0.42),
         "為什麼主路由不是 270M，是它", size=16, bold=True, color=TEALDK)
_pts = [
    ("🧭", "先分類、再生成", "每句先由 FastText 判斷「做什麼」，高信心直接跳過 LLM——"
                          "實測過半句子根本不勞 270M 出手，一句回答快在這裡"),
    ("⚖️", "確定性路由", "向量比對、同句永遠同答案，沒有 LLM 的浮點抖動——"
                       "雙平台（WIN11 / RPi5）行為一致的基石"),
    ("訓", "從零訓練", "FastText 出廠不帶任何類別（它是演算法不是現成模型）；"
                      "14 類依倉管工具集自定義，語料與 LLM 微調同源、判斷口徑一致"),
]
for i, (ic, h, d) in enumerate(_pts):
    y = _ly + Inches(0.78) + Inches(0.94) * i
    dot_icon(s, MX + Inches(0.35), y + Inches(0.02), ic, circle=TEAL, gcolor=WHITE)
    add_text(s, MX + Inches(0.95), y - Inches(0.06), Inches(5.3), Inches(0.34),
             h, size=13.5, bold=True, color=DARK)
    add_text(s, MX + Inches(0.95), y + Inches(0.26), Inches(5.35), Inches(0.62),
             d, size=12, color=GREY44)

# 右卡：14 類一覽（標籤名=工具名）
_rx = MX + Inches(6.85)
_rw = SLIDE_W - _rx - MX
add_round(s, _rx, _ly, _rw, Inches(3.72), fill=TEALBG)
add_text(s, _rx + Inches(0.32), _ly + Inches(0.22), _rw - Inches(0.64), Inches(0.4),
         "14 類意圖一覽（標籤名＝工具名）", size=15, bold=True, color=TEALDK)
_groups = [
    ("查詢 × 8", "庫存 / 進出 / 缺貨 / 熱銷 / 連帶 / 比倉 / 比期間 / 清單"),
    ("鑑識 × 1", "查帳追根因（search_log 自動比對進出差異）"),
    ("控制 × 5", "改設定 / 警示 / 採購單 / 報表 / 白名單腳本"),
]
for i, (g, items) in enumerate(_groups):
    y = _ly + Inches(0.78) + Inches(0.84) * i
    add_text(s, _rx + Inches(0.32), y, Inches(1.5), Inches(0.34),
             g, size=13, bold=True, color=DARK)
    add_text(s, _rx + Inches(0.32), y + Inches(0.32), _rw - Inches(0.64), Inches(0.5),
             items, size=12, color=GREY44)
add_text(s, _rx + Inches(0.32), _ly + Inches(3.28), _rw - Inches(0.64), Inches(0.36),
         "分類結果直接就是要呼叫的功能——沒有中間翻譯層", size=12, bold=True, color=TEALDK)
set_notes(s, "FastText 是 Facebook 開源的輕量文字分類技術：詞向量平均＋線性分類器，"
             "沒有生成、沒有推理，所以毫秒級、確定性（同句同答案）。重點澄清：FastText "
             "出廠不帶任何分類類別——官方只有語言辨識模型和詞向量，意圖分類一定要自備語料"
             "從零訓練；我們的 14 類直接對應倉管工具名，訓練語料跟 LLM 微調同源（5,849 筆），"
             "所以 clf 和 LLM 判斷口徑一致。分工哲學：clf 管「做什麼」（路由），270M 管"
             "「參數是什麼」（抽取）——路由要穩定用分類器，抽取要泛化用 LLM，各用其長。"
             "高信心（conf≥0.8）且不需參數的功能直接跳過 LLM，是「一秒回答」的主要來源。"
             "彩蛋（可講）：這顆分類器曾因 fasttext×numpy2 不相容在展示機上靜默死亡三週，"
             "期間 100+ 輪、六套回歸照樣雙平台 100%——LLM+校正層獨自扛住，正是縱深防禦的實證；"
             "修復後展示機過半問題毫秒級回答、路由徽章首次亮起。")
pn(s)


# ─── S6 為什麼是 270M（尺寸定位）───────────────────────────
s = slide_blank()
title_bar(s, "WHY SO SMALL", "270M：故意選小，不是能力不足")
kpi_row(s, Inches(1.75), [
    ("270M", "本專案模型參數"),
    ("3-8B", "業界主流 Agent"),
    ("20–30 t/s", "RPi5 CPU 推論速度"),
    ("0 元", "雲端 API 費用"),
], num_size=27)
add_round(s, MX, Inches(3.35), Inches(11.87), Inches(3.3), fill=LIGHT, shadow=True)
add_text(s, MX + Inches(0.4), Inches(3.6), Inches(11.0), Inches(0.45),
         "小模型的三個好處", size=17, bold=True, color=TEALDK)
pts = [
    ("💰", "成本極低", "跑在一台樹莓派 CPU 上，不需 GPU、不付雲端 API 費用"),
    ("🔌", "離線可用", "展場用手機熱點就能跑，不依賴網路，資料不出場"),
    ("🎯", "定位精準", "270M 當「路由器」而非「決策者」——分類意圖它綽綽有餘，"
                    "精確計算交給規則層，反而更穩"),
]
for i, (ic, h, d) in enumerate(pts):
    y = Inches(4.15) + Inches(0.78) * i
    dot_icon(s, MX + Inches(0.4), y + Inches(0.02), ic, circle=TEAL, gcolor=WHITE)
    add_text(s, MX + Inches(1.0), y + Inches(0.09), Inches(1.7), Inches(0.5),
             h, size=14.5, bold=True, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, MX + Inches(2.7), y + Inches(0.09), Inches(8.5), Inches(0.6),
             d, size=13, color=GREY44, anchor=MSO_ANCHOR.MIDDLE)
set_notes(s, "回應老闆可能的疑問：為什麼不用更大的模型？答案是刻意的取捨。"
             "270M 讓整套系統能在便宜硬體上離線跑，而把精確度靠規則層補足——"
             "這正是這個展示的技術亮點：用最小的模型做到可用的品質。")
pn(s)

print("S5-S6 done")


# ─── S7 測試哲學（章節轉場，深底）─────────────────────────
s = slide_blank()
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=DARK)
add_rect(s, 0, 0, Inches(0.14), SLIDE_H, fill=TEAL)
add_text(s, MX, Inches(2.4), Inches(11.8), Inches(0.4),
         "TESTING PHILOSOPHY", font=FONT_EN, size=13, bold=True, color=TEAL)
add_text(s, MX, Inches(2.9), Inches(11.8), Inches(1.4),
         "展場訪客會亂打字，\n所以我們先自己把它打壞", size=34, bold=True, color=WHITE,
         line_spacing=1.1)
add_text(s, MX, Inches(4.6), Inches(11.0), Inches(0.9),
         "小模型的先天限制是輸出不穩定。與其祈禱它不出錯，不如用大量真實情境反覆攻擊，"
         "把每個破口修掉、再存成回歸測試，確保永不重犯。", size=15, color=GREYBB,
         line_spacing=1.35)
set_notes(s, "測試段開場。核心理念：不相信小模型天生可靠，而是靠測試把品質「逼」出來。"
             "展場訪客會用各種亂七八糟的方式打字，我們的做法是自己先模擬這些情境攻擊系統。")
pn(s)

# ─── S8 六套回歸測試（總覽）──────────────────────────────
s = slide_blank()
title_bar(s, "SIX TEST SUITES", "六套回歸測試，全部雙平台 100%")
suites = [
    ("守衛庫", f"{GUARD_ZH} 句", "每次修 bug 都存成守衛，防止回退", "🛡️"),
    ("短句全枚舉", "953 句", "60 商品 × 模板全展開，證明短句 100%", "🔬"),
    ("多輪對話全枚舉", "1980 情境", "首句 + 追問，證明「記得上一個商品」", "💬"),
    ("未知商品抗性", "60+ 情境", "新增商品後查、問不存在的、撞名、亂取名", "🧩"),
    ("81 題標準集", "99%", "路由準確率基準測試", "🎯"),
    ("OOV 口語集", "98%", "沒見過的口語 / 錯字容錯", "🗣️"),
]
ICON["🧩"] = "測"
y0 = Inches(1.6); rh = Inches(0.86); rgap = Inches(0.13)
for i, (name, num, desc, ic) in enumerate(suites):
    y = y0 + (rh + rgap) * i
    add_round(s, MX, y, Inches(11.87), rh, fill=LIGHT, shadow=True)
    add_icon_circle(s, MX + Inches(0.22), y + Inches(0.17), Inches(0.52), ic,
                    circle=TEAL, gsize=15)
    add_text(s, MX + Inches(1.0), y + Inches(0.05), Inches(3.0), Inches(0.76),
             name, size=15.5, bold=True, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, MX + Inches(3.9), y + Inches(0.05), Inches(1.9), Inches(0.76),
             num, font=FONT_EN, size=18, bold=True, color=TEALDK,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, MX + Inches(5.9), y + Inches(0.05), Inches(5.8), Inches(0.76),
             desc, size=12.5, color=GREY44, anchor=MSO_ANCHOR.MIDDLE)
set_notes(s, "六套測試各有分工。守衛庫防回退、短句全枚舉證明產品本體、多輪對話證明"
             "上下文記憶、標準集與 OOV 集是路由準確率基準。全部在 Windows 開發機與"
             "樹莓派兩個平台都跑到 100% 才算過。")
pn(s)

print("S7-S8 done")


# ─── S9 全枚舉方法（把「多輪」變成可數的）──────────────────
s = slide_blank()
title_bar(s, "EXHAUSTIVE SWEEP", "怎麼證明「多輪對話」不會出錯？把它數出來")
add_text(s, MX, Inches(1.45), Inches(11.8), Inches(0.7),
         "隨機測試永遠測不完，只是「抽樣」。真正的保證來自「窮舉」——"
         "把整個空間定義成可以數的組合，全部跑一遍。", size=14, color=GREY55,
         line_spacing=1.3)
# 公式卡
add_round(s, MX, Inches(2.4), Inches(11.87), Inches(1.5), fill=DARK, shadow=True)
add_text(s, MX + Inches(0.5), Inches(2.62), Inches(11.0), Inches(0.4),
         "多輪短句空間 = 可窮舉", size=14, bold=True, color=TEAL)
add_rich(s, MX + Inches(0.5), Inches(3.05), Inches(11.0), Inches(0.7),
         [{"text": "60 ", "font": FONT_EN, "size": 26, "bold": True, "color": WHITE},
          {"text": "商品   ×   ", "size": 18, "color": GREYBB},
          {"text": "33 ", "font": FONT_EN, "size": 26, "bold": True, "color": WHITE},
          {"text": "種追問形   =   ", "size": 18, "color": GREYBB},
          {"text": "1980 ", "font": FONT_EN, "size": 26, "bold": True, "color": TEAL},
          {"text": "個兩輪對話", "size": 18, "color": WHITE}],
         anchor=MSO_ANCHOR.MIDDLE)
# 追問形六族
add_text(s, MX, Inches(4.25), Inches(11.0), Inches(0.4),
         "33 種追問形 = 從實測 bug 反推的六大族", size=15, bold=True, color=DARK)
fams = [("代詞", "那個呢 / 它還剩幾個"), ("功能詞", "進出 / 到期 / 多少錢"),
        ("倉別", "南倉呢 / 北倉多少 / 南"), ("語助詞", "呢 / 咧"),
        ("寫入", "北倉進 20 個"), ("錯字注音", "那個近出紀錄呢 / 安全ㄎ存")]
for i, (h, ex) in enumerate(fams):
    r, c = divmod(i, 3)
    x = MX + (Inches(3.9) + Inches(0.08)) * c
    y = Inches(4.8) + Inches(0.82) * r
    add_round(s, x, y, Inches(3.9), Inches(0.68), fill=TEALBG)
    add_text(s, x + Inches(0.2), y + Inches(0.09), Inches(1.2), Inches(0.5),
             h, size=13.5, bold=True, color=TEALDK, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(1.35), y + Inches(0.09), Inches(2.4), Inches(0.5),
             ex, size=12, color=GREY55, anchor=MSO_ANCHOR.MIDDLE)
set_notes(s, "全枚舉的核心價值：把「多輪」這個看似無邊無際的東西，定義成一個可以數的空間——"
             "60 個商品乘上 33 種追問形，就是 1980 個兩輪對話。33 種追問形不是隨便想的，"
             "是從前面四輪隨機測試挖到的真 bug 反推出來的六大族。")
pn(s)

# ─── S10 全枚舉成果（前後對比）────────────────────────────
s = slide_blank()
title_bar(s, "SWEEP RESULTS", "全枚舉抓出隨機測試漏掉的破口，一次修光")
# 對比兩張大卡
add_round(s, MX, Inches(1.85), Inches(5.75), Inches(2.3), fill=LIGHT, shadow=True)
add_text(s, MX + Inches(0.4), Inches(2.1), Inches(5.0), Inches(0.4),
         "隨機採樣（4 輪）", size=15, bold=True, color=GREY55)
add_text(s, MX + Inches(0.4), Inches(2.55), Inches(5.0), Inches(0.9),
         "27", font=FONT_EN, size=54, bold=True, color=GREY77)
add_text(s, MX + Inches(2.4), Inches(2.95), Inches(3.0), Inches(0.5),
         "個 bug / 102 情境", size=14, color=GREY55)
add_text(s, MX + Inches(0.4), Inches(3.6), Inches(5.0), Inches(0.4),
         "只看到冰山一角", size=13, color=GREY77)

rx = Inches(6.85)
add_round(s, rx, Inches(1.85), Inches(5.75), Inches(2.3), fill=TEALBG, shadow=True)
add_text(s, rx + Inches(0.4), Inches(2.1), Inches(5.0), Inches(0.4),
         "全枚舉（1 次）", size=15, bold=True, color=TEALDK)
add_text(s, rx + Inches(0.4), Inches(2.55), Inches(5.0), Inches(0.9),
         "29", font=FONT_EN, size=54, bold=True, color=TEALDK)
add_text(s, rx + Inches(2.4), Inches(2.95), Inches(3.0), Inches(0.5),
         "個 bug / 1980 情境", size=14, color=TEALDK)
add_text(s, rx + Inches(0.4), Inches(3.6), Inches(5.0), Inches(0.4),
         "通過率 98.5%，破口全面現形", size=13, color=TEALDK)
# 根因總結
add_round(s, MX, Inches(4.45), Inches(11.87), Inches(2.1), fill=DARK, shadow=True)
add_text(s, MX + Inches(0.4), Inches(4.68), Inches(11.0), Inches(0.4),
         "29 個破口 → 兩個根因 → 兩處修補", size=15, bold=True, color=TEAL)
add_rich(s, MX + Inches(0.4), Inches(5.2), Inches(11.2), Inches(0.55),
         [{"text": "28 個  ", "size": 15, "bold": True, "color": WHITE},
          {"text": "共用字的商品（USB風扇 / 素T…）進出查詢把別的品項也算進去 → 加分數門檻，一修全清",
           "size": 13.5, "color": GREYBB}])
add_rich(s, MX + Inches(0.4), Inches(5.78), Inches(11.2), Inches(0.55),
         [{"text": "1 個   ", "size": 15, "bold": True, "color": WHITE},
          {"text": "「清潔手套」的「清潔」被誤當類別詞 → 商品名優先，讓給完整比對",
           "size": 13.5, "color": GREYBB}])
add_text(s, MX + Inches(0.4), Inches(6.28), Inches(11.0), Inches(0.3),
         "關鍵發現：這些其實是「單句就會錯」的 bug，只是舊測試沒涵蓋到這幾個商品的這個功能組合",
         size=12, color=GREY77)
set_notes(s, "全枚舉的成果。左右對比：四輪隨機測試挖到 27 個，但全枚舉一次照出 29 個——"
             "而且看得到全貌（1980 情境）。這 29 個高度集中在兩個根因，兩處修補就解決。"
             "最重要的發現：這些其實是單句就存在的 bug，全枚舉把舊測試的盲區補上了。")
pn(s)

print("S9-S10 done")


# ─── S11 收斂軌跡（原生柱狀圖）────────────────────────────
s = slide_blank()
title_bar(s, "CONVERGENCE", f"100+ 輪迭代，守衛庫從 138 句長到 {GUARD_ZH} 句")
# 左：守衛庫成長 折線/柱
chart_data = CategoryChartData()
chart_data.categories = ["初期", "conv100\n收斂", "第二戰役\n短句認證", "多輪\n全枚舉", "語音 POC\n真人聲"]
chart_data.add_series("守衛庫句數", (138, 352, 855, 978, GUARD_ZH))
gx, gy = Inches(0.73), Inches(1.75)
gcx = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, gx, gy,
                         Inches(6.1), Inches(3.6), chart_data).chart
gcx.has_legend = False
gcx.has_title = True
gcx.chart_title.text_frame.text = "守衛庫句數成長"
for r in gcx.chart_title.text_frame.paragraphs[0].runs:
    r.font.size = Pt(13); r.font.name = FONT_ZH; r.font.bold = True; r.font.color.rgb = DARK
plot = gcx.plots[0]
plot.has_data_labels = True
plot.data_labels.font.size = Pt(12)
plot.data_labels.font.name = FONT_EN
plot.data_labels.number_format = "0"
plot.data_labels.number_format_is_linked = False
plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
ser = plot.series[0]
ser.format.fill.solid(); ser.format.fill.fore_color.rgb = TEAL
cat_ax = gcx.category_axis
cat_ax.tick_labels.font.size = Pt(11); cat_ax.tick_labels.font.name = FONT_ZH
val_ax = gcx.value_axis
val_ax.tick_labels.font.size = Pt(10); val_ax.tick_labels.font.name = FONT_EN
val_ax.has_major_gridlines = False
val_ax.maximum_scale = 1000

# 右：多輪 bug 軌跡 + 說明
add_round(s, Inches(7.1), Inches(1.75), Inches(5.5), Inches(3.6), fill=LIGHT, shadow=True)
add_text(s, Inches(7.4), Inches(1.98), Inches(5.0), Inches(0.4),
         "多輪真 bug 軌跡", size=15, bold=True, color=DARK)
add_rich(s, Inches(7.4), Inches(2.5), Inches(5.0), Inches(0.5),
         [{"text": "r32→r35：", "size": 14, "color": GREY55},
          {"text": "  6 → 7 → 6 → 8", "font": FONT_EN, "size": 20, "bold": True, "color": AMBER}])
add_text(s, Inches(7.4), Inches(3.15), Inches(5.0), Inches(1.0),
         "四輪隨機測試沒收斂——每輪只看到不同角落。\n"
         "這反而證明：光靠隨機測不夠，必須全枚舉。", size=13, color=GREY44, line_spacing=1.35)
add_rich(s, Inches(7.4), Inches(4.35), Inches(5.0), Inches(0.7),
         [{"text": "全枚舉修完：", "size": 14, "color": GREY55},
          {"text": "  1980 / 1980 全綠", "size": 15, "bold": True,
           "color": TEALDK}])
# 底部一句話
add_round(s, MX, Inches(5.7), Inches(11.87), Inches(0.95), fill=TEALBG)
add_text(s, MX + Inches(0.4), Inches(5.9), Inches(11.0), Inches(0.6),
         f"每修一個 bug 就存一句守衛 → {GUARD_ZH} 句形成防護網，任何改動只要跑回歸就知道有沒有踩壞舊功能",
         size=13.5, color=TEALDK, anchor=MSO_ANCHOR.MIDDLE)
set_notes(s, "收斂的量化證據。左邊柱狀圖：守衛庫從 138 句一路長到 866 句，每個數字都是"
             "累積的防護網。右邊：多輪 bug 軌跡 6-7-6-8 沒收斂，正好證明隨機測試的極限，"
             "催生了全枚舉的方法。")
pn(s)

# ─── S11b 自動測試迭代流程（收尾階段報告）★ 2026-08-06 ──────────
s = slide_blank()
title_bar(s, "AUTOMATED TESTING",
          "自動測試迭代：百句實測 → 修 → 回歸；逐條功能線收斂，換角度再測")

# 上半：五步循環流程卡
_steps = [
    ("🎲", "產生測試批",   "每輪 100 句全新、\n刻意換角度：口語、\n邊界、劇情、語音錯字"),
    ("🖥️", "瀏覽器自動實測", "程式模擬真訪客\n打字送出，逐句\n截圖存證"),
    ("🔍", "逐句判定",     "✅ 正確\n🟡 合理降級\n❌ 破口"),
    ("🛠️", "追根因修復",   "看執行 log 定位\n卡在哪一層，\n對症下藥不亂猜"),
    ("🛡️", "回歸與部署", f"中 {GUARD_ZH}＋英 {GUARD_EN} 句\n守衛防「修A壞B」，\n三處版控雙機同步"),
]
_sw, _sh, _gap = Inches(2.24), Inches(2.05), Inches(0.19)
_sx0, _sy0 = MX, Inches(1.62)
for _i, (_g, _t, _d) in enumerate(_steps):
    _sx = _sx0 + _i * (_sw + _gap)
    add_round(s, _sx, _sy0, _sw, _sh, fill=LIGHT, shadow=True)
    add_icon_circle(s, _sx + Inches(0.12), _sy0 + Inches(0.14), 0.44, _g,
                    circle=TEAL, gsize=15)
    add_text(s, _sx + Inches(0.62), _sy0 + Inches(0.17), _sw - Inches(0.7), Inches(0.4),
             f"{_i+1}. {_t}", size=13.5, bold=True, color=DARK)
    add_text(s, _sx + Inches(0.16), _sy0 + Inches(0.7), _sw - Inches(0.3), Inches(1.25),
             _d, size=12, color=GREY44, line_spacing=1.25)
    if _i < 4:
        add_arrow(s, _sx + _sw - Inches(0.02), _sy0 + Inches(0.85), Inches(0.23), Inches(0.3))
# 回繞說明（第五卡下方 → 第一卡）
add_text(s, _sx0, _sy0 + _sh + Inches(0.08), Inches(11.87), Inches(0.34),
         "↺  一輪修完不算完——換一批全新句子、換一種訪客講話方式，再跑下一輪，直到連續兩輪幾乎抓不到新問題",
         size=12.5, color=TEALDK, bold=True)

# 下半左：破口收斂柱狀圖（原生 chart）
# 2026-08-06 user 校正：原本只畫英文互動五輪（38→3）＝看起來全案收斂完，
#   但那只是**一條功能線**。之後換新角度（排程專項）又抓到 4 個、實際
#   使用回報 4 件 ⇒ 補上第六柱，讓「換角度就會再冒」這件事在圖上看得見，
#   老闆才不會誤以為已結案。
_cd = CategoryChartData()
_cd.categories = ["第一輪", "第二輪", "第三輪", "第四輪", "第五輪", "換角度\n(排程專項)"]
_cd.add_series("每輪新抓到的問題數", (38, 31, 19, 9, 3, 4))
_gc = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, MX, Inches(4.42),
                         Inches(5.9), Inches(2.1), _cd).chart
_gc.has_legend = False
_gc.has_title = True
_gc.chart_title.text_frame.text = "破口收斂軌跡：換新角度仍會再冒"
for _r in _gc.chart_title.text_frame.paragraphs[0].runs:
    _r.font.size = Pt(12.5); _r.font.name = FONT_ZH; _r.font.bold = True; _r.font.color.rgb = DARK
_pl = _gc.plots[0]
_pl.has_data_labels = True
_pl.data_labels.font.size = Pt(12)
_pl.data_labels.font.name = FONT_EN
_pl.data_labels.number_format = "0"
_pl.data_labels.number_format_is_linked = False
_pl.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
_pl.series[0].format.fill.solid()
_pl.series[0].format.fill.fore_color.rgb = TEAL
_gc.category_axis.tick_labels.font.size = Pt(11)
_gc.category_axis.tick_labels.font.name = FONT_ZH
_gc.value_axis.tick_labels.font.size = Pt(10)
_gc.value_axis.tick_labels.font.name = FONT_EN
_gc.value_axis.has_major_gridlines = False

# 下半右：收斂判準卡
add_round(s, Inches(6.95), Inches(4.42), Inches(5.65), Inches(2.1), fill=LIGHT, shadow=True)
add_text(s, Inches(7.22), Inches(4.58), Inches(5.1), Inches(0.38),
         "收斂判準（工程標準）", size=14, bold=True, color=DARK)
add_rich(s, Inches(7.22), Inches(5.02), Inches(5.2), Inches(0.42),
         [{"text": "連續兩輪新問題 ≤ 5 且全為優雅降級  ", "size": 12.5, "color": GREY44},
          {"text": "→ 英文互動線已達標", "size": 13, "bold": True, "color": TEALDK}])
add_text(s, Inches(7.22), Inches(5.5), Inches(5.2), Inches(0.95),
         "但收斂是「逐條功能線」達標，不是全案結案：換新角度測（排程專項）\n"
         "仍抓到 4 個破口，實際使用回報 4 件。⇒ 下週持續：新角度輪 + 守衛\n"
         "庫補課（近期修復尚未全數納入回歸語料）。",
         size=11.5, color=GREY44, line_spacing=1.28)

# 底部狀態條（交付訊息）
add_round(s, MX, Inches(6.72), Inches(11.87), Inches(0.6), fill=TEALBG)
add_rich(s, MX + Inches(0.35), Inches(6.83), Inches(11.3), Inches(0.4),
         [{"text": "收尾階段   ", "size": 13.5, "bold": True, "color": TEALDK},
          {"text": "▸  下週續辦：新角度輪、守衛庫補課、真人語音驗證   ",
           "size": 13, "color": GREY44},
          {"text": "▸  展前（9/2）完成交付", "size": 13.5, "bold": True, "color": DARK}])
set_notes(s, "這頁講測試方法論與現況。核心觀念：不是測一次就好，而是「百句實測→修→"
             "全量回歸」的循環，每輪刻意換訪客講話的角度（口語、邊界、連續對話、語音"
             "錯字），修完一輪再跑下一輪。\n"
             "柱狀圖：英文互動線五輪從 38 個問題降到 3 個，連續兩輪低於 5 ⇒ **那條線**"
             "達到收斂判準。但第六柱是關鍵——換一個沒測過的角度（排程專項百句）又抓到 "
             "4 個，同期實際使用回報 4 件。\n"
             "★ 這頁要傳達的重點不是『測完了』，而是『收斂是逐條功能線達標，換角度就會"
             "再冒』。老闆若問「那到底何時算好」：判準是每條功能線各自連兩輪 ≤5，"
             "加上實際使用回報趨近於零；目前英文互動線已達標，排程線剛開始測。\n"
             f"守衛庫補課：ZH 1122→{GUARD_ZH}、EN 892→{GUARD_EN}（7/20、7/25 後停止成長的缺口已補），"
             "下週要補課，否則新修的功能沒有防退步保護。\n"
             "結論：收尾階段，下週續辦新角度輪＋守衛補課＋真人語音驗證，展前（9/2）完成交付。")
pn(s)

# ─── S11c 測過哪些角度／還剩哪些（2026-08-06 user 提問：會被追問）──────
s = slide_blank()
title_bar(s, "TEST COVERAGE", "「換角度」換的是什麼：已測 35+ 輪的角度盤點與剩餘缺口")
add_text(s, MX, Inches(1.32), Inches(11.8), Inches(0.4),
         "每輪刻意換一種「訪客會怎麼講話」的角度，而不是同一批句子重跑。"
         "下表是已覆蓋的角度、以及展前還要補的三塊。",
         size=13, color=GREY55)

_cov = [
    ("訪客講話方式", "口語簡稱、模糊描述、禮貌繞圈、一句多意圖、代稱追問", True),
    ("輸入型態", "打字錯字、注音殘留、黏字漏空格、大小寫、標點缺失", True),
    ("對話結構", "多輪劇情、上下文接續、反悔改單、中途插話、確認代按", True),
    ("業務語境", "營運詞彙（滯銷/呆料/撐天）、期間表達、跨倉比較、排除語境", True),
    ("邊界與異常", "查無商品、數量為零/負、超量調撥、不支援期間、重複排程", True),
    ("搗蛋與離題", "閒聊、辱罵、注入攻擊（drop all tables）、問系統身世", True),
    ("語音輸入", "ASR 錯字型態、同音字、四種腔調 TTS、三層噪音", True),
    ("功能面全覆蓋", "七大查詢 × 三類 Agent 工具 × 排程/警示全生命週期", True),
    ("短句全枚舉", "1-4 字的所有合理組合（把「多輪」變成可數空間）", True),
]
_gap = [
    ("① 真人語音批 · 100 句", "新麥克風到貨後重錄；重點打「換人、換口音」的新錯法",
     "麥克風到貨當天"),
    ("② 展場情境批 · 100 句", "站著問的短促語氣、被打斷重問、旁人插話、邊看螢幕邊改口",
     "下週一～二"),
    ("③ 併發壓力批 · 5 路 ×20", "排隊輪流講話、麥克風連續佔用、同時送出的搶鎖情境",
     "下週三"),
]
_y = Inches(1.86)
add_round(s, MX, _y, Inches(7.15), Inches(4.42), fill=LIGHT, shadow=True)
add_text(s, MX + Inches(0.26), _y + Inches(0.14), Inches(6.6), Inches(0.34),
         "✅ 已覆蓋的角度（35+ 輪累積）", size=13.5, bold=True, color=TEALDK)
_ry = _y + Inches(0.56)
for _nm, _dsc, _ in _cov:
    add_text(s, MX + Inches(0.26), _ry, Inches(2.35), Inches(0.3),
             f"· {_nm}", size=11, bold=True, color=DARK)
    add_text(s, MX + Inches(2.62), _ry, Inches(4.4), Inches(0.3),
             _dsc, size=10, color=GREY55)
    _ry += Inches(0.42)

add_round(s, Inches(8.18), _y, Inches(4.42), Inches(4.42), fill=WHITE,
          line=CORAL, shadow=True)
add_text(s, Inches(8.44), _y + Inches(0.14), Inches(3.9), Inches(0.34),
         "📋 下週排程：三個新角度", size=13.5, bold=True, color=CORAL)
_gy = _y + Inches(0.6)
for _nm, _why, _when in _gap:
    add_text(s, Inches(8.44), _gy, Inches(3.9), Inches(0.3),
             f"· {_nm}", size=11.5, bold=True, color=DARK)
    add_text(s, Inches(8.44), _gy + Inches(0.3), Inches(3.9), Inches(0.62),
             _why, size=10, color=GREY55, line_spacing=1.18)
    add_text(s, Inches(8.44), _gy + Inches(0.94), Inches(3.9), Inches(0.28),
             f"時程：{_when}", size=10, bold=True, color=TEALDK)
    _gy += Inches(1.32)

add_round(s, MX, Inches(6.48), Inches(11.87), Inches(0.72), fill=DARK, shadow=True)
add_rich(s, MX + Inches(0.35), Inches(6.6), Inches(11.2), Inches(0.5),
         [[{"text": "下週產出  ", "size": 12.5, "bold": True, "color": TEAL},
           {"text": "三批共 220 句實測 → 修復 → 全量回歸 → 納入守衛庫",
            "size": 12, "bold": True, "color": WHITE},
           {"text": "——角度是先宣告再打，不是測完才回頭找名目；"
                    "每批的破口數與修復清單都會留紀錄，收斂判準同前頁（連兩輪 ≤5 且全為優雅降級）。",
            "size": 11.5, "color": GREYBB}]],
         anchor=MSO_ANCHOR.MIDDLE)
set_notes(s, "★這頁回答「換角度是換什麼、還有哪些沒測」——老闆一定追問，而且要看到計畫"
             "不是「再測測看」。\n"
             "左欄九個角度是 35+ 輪累積的實際批次。可以挑兩個講：短句全枚舉把 1-4 字的"
             "所有合理組合窮舉，等於把「多輪對話」這個看似無邊際的東西變成可數空間；"
             "搗蛋批連 SQL 注入（drop all tables）都測過。\n"
             "右欄是**下週的具體排程**，每批都講得出打什麼句子：\n"
             "① 真人語音 100 句——現有語音數據全來自單一錄音者加合成音，換人換麥克風"
             "必然冒新錯法（今天用四種 TTS 腔調挖過，800 句挖到 146 句未命中）。"
             "麥克風到貨當天就錄。\n"
             "② 展場情境 100 句——這是全新角度：訪客是**站著**問的，語氣短促、"
             "容易被旁人打斷、會邊看螢幕邊改口。這些句型跟坐著打字完全不同。\n"
             "③ 併發壓力 5 路×20——已測過 5 路同時送出不崩，但沒測過「排隊輪流講話」"
             "「麥克風被連續佔用」這種展場實況。\n"
             "★ 若老闆問「這樣要測到什麼時候」：判準明確（連兩輪新破口 ≤5 且全為優雅降級），"
             "英文互動線已達標，這三批是最後的覆蓋缺口，展前完成。\n"
             "★ 若老闆質疑在空轉：反過來看——能列出「還沒測什麼」而且說得出每批要打什麼、"
             "哪天做，代表測試是有地圖的。怕的不是有缺口，是不知道缺在哪。")
pn(s)

# ─── S14a 英文版 · 為何不是翻譯（路線決策）★ ──────────────────────
s = slide_blank()
title_bar(s, "ENGLISH BUILD · 路線", "做英文版不是翻譯——翻譯會讓招牌能力全滅")
add_text(s, MX, Inches(1.35), Inches(11.8), Inches(0.4),
         "老闆要全英文版。動工前先用探針餵分級英文句給現有模型，量出「翻譯路線」到底會壞在哪——結論決定了整條路線。",
         size=13, color=GREY55)
# 左：探針結果（翻譯路線會壞的地方）
add_round(s, MX, Inches(1.95), Inches(5.82), Inches(3.5), fill=LIGHT, shadow=True)
add_text(s, MX + Inches(0.3), Inches(2.1), Inches(5.2), Inches(0.36),
         "探針實測：純翻譯會壞在哪", size=14, bold=True, color=DARK)
probe = [
    ("✓", "乾淨查詢", "how many bluetooth earphones left → 答對", TEALDK),
    ("×", "英文錯字", "earphon → 完全對不到", CORAL),
    ("×", "模糊描述", "the thing that charges phone → 不懂", CORAL),
    ("×", "寫入 / 調貨 / RCA", "add / move / why → 全歸零查詢", CORAL),
]
for i, (mk, tag, ex, col) in enumerate(probe):
    y = Inches(2.55) + Inches(0.68) * i
    add_text(s, MX + Inches(0.3), y, Inches(0.3), Inches(0.3), mk,
             font=FONT_EN, size=14, bold=True, color=col)
    add_text(s, MX + Inches(0.66), y - Inches(0.02), Inches(1.9), Inches(0.32),
             tag, size=12.5, bold=True, color=col)
    add_text(s, MX + Inches(0.66), y + Inches(0.28), Inches(4.8), Inches(0.32),
             ex, size=11, color=GREY55)
# 右：三條路線比較
add_round(s, MX + Inches(6.05), Inches(1.95), Inches(5.82), Inches(3.5), fill=LIGHT, shadow=True)
add_text(s, MX + Inches(6.35), Inches(2.1), Inches(5.2), Inches(0.36),
         "三條路線，選了中間那條", size=14, bold=True, color=DARK)
routes = [
    ("翻譯層", "最省事，但容錯層先崩就輪不到後面——等於白費版", GREY77, False),
    ("補英文語料微調", "Gemma 英文底子還在，教它這套系統的 tool 慣例", TEALDK, True),
    ("全部重訓", "成本最高，但 base 英文能力本來就在，沒必要", GREY77, False),
]
for i, (nm, desc, col, pick) in enumerate(routes):
    y = Inches(2.6) + Inches(0.92) * i
    add_round(s, MX + Inches(6.35), y, Inches(5.22), Inches(0.78),
              fill=(TEALBG if pick else WHITE), line=(TEAL if pick else GREYE6),
              line_w=(1.3 if pick else 0.6))
    add_text(s, MX + Inches(6.55), y + Inches(0.06), Inches(1.85), Inches(0.32),
             ("★ " if pick else "") + nm, size=12.5, bold=True, color=col)
    add_text(s, MX + Inches(6.55), y + Inches(0.38), Inches(4.85), Inches(0.34),
             desc, size=10.5, color=GREY55)
# 底部：微調效益數字
add_round(s, MX, Inches(5.62), Inches(11.87), Inches(1.54), fill=DARK, shadow=True)
add_text(s, MX + Inches(0.4), Inches(5.78), Inches(11.1), Inches(0.34),
         "三方對照：同一份 34 句英文評測集（本機 llama.cpp 實跑）",
         size=12.5, bold=True, color=TEAL)
tri = [("base 未微調", "11%", "看得懂英文，但不知道該叫哪個 tool"),
       ("中文微調版", "32%", "tool 慣例**跨語言遷移**——用中文學的可套到英文"),
       ("英文微調版", "73%", "基本查詢 12/12、錯字全中、RCA 3/3")]
for i, (nm, sc, note) in enumerate(tri):
    x = MX + Inches(0.4) + Inches(3.85) * i
    add_text(s, x, Inches(6.18), Inches(1.5), Inches(0.42), sc,
             font=FONT_EN, size=22, bold=True,
             color=(TEAL if i == 2 else (WHITE if i == 1 else GREY77)))
    add_text(s, x + Inches(1.45), Inches(6.2), Inches(2.3), Inches(0.36),
             nm, size=11.5, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x, Inches(6.66), Inches(3.6), Inches(0.34),
             note.replace("**", ""), size=10.5, color=GREYBB)
set_notes(s, "★英文版路線決策頁。老闆要全英文版，我沒有直接開始翻譯——先用探針餵分級英文句"
             "給現有的中文微調模型，量出「純翻譯路線」會壞在哪。結果很清楚：乾淨查詢答得對"
             "（Gemma 的英文底子還在），但**招牌能力全滅**——英文錯字、模糊描述、寫入/調貨/RCA "
             "意圖全部救不到。而容錯層正是這個 demo 的賣點，翻譯層先崩就輪不到後面，等於做出一個"
             "「白費版」。所以路線選**補英文語料微調**：不必全部重訓（base 英文能力本來就在），"
             "但要教它這套系統的 tool 慣例。右下三方對照是實測數字：base 未微調只有 11%（看得懂"
             "英文、不知道叫哪個 tool）；有趣的是**中文微調版跑英文有 32%**——證明 tool 慣例會"
             "跨語言遷移，用中文學的可以套到英文，這是微調買到的『領域判斷』而不只是語言；"
             "英文專訓版 73%。這頁的重點是：**決策有數據支撐，不是拍腦袋選路線**。")
pn(s)


# ─── S14b 英文版 · 19 類移植坑（真正的工作量）★ ────────────────────
s = slide_blank()
title_bar(s, "ENGLISH BUILD · 真正的工作量", "移植的難處不在模型，在散落各處的語言假設")
add_text(s, MX, Inches(1.35), Inches(11.8), Inches(0.4),
         "模型換好只是開始。真正吃時間的是三個月為中文磨出來的規則層——每一條都藏著「這是中文」的隱含假設。",
         size=13, color=GREY55)
traps = [
    ("詞表是中文", "`_ALL_INTENT_WORDS` 等 21 個詞表全中文 → 英文句一個都不命中",
     "英文句被判成「只有商品名沒動作」，整批轉 clarify", NAVY),
    ("對照表的**鍵**是中文", "`{\"電子\": \"electronics\"}` ——值是英文、鍵是中文",
     "6 個類別 5 個查不到，整條類別查詢功能靜默失效", TEALDK),
    ("門檻用中文字元數", "長句判定 >30 字元；英文字元數是中文的 2-3 倍",
     "正常英文句被當長句，幾乎全部繞過 LLM", TEAL),
    ("演算法假設中文形態", "`split()[0]` 剝規格尾巴——中文沒空白所以安全",
     "`Wireless Mouse` 被腰斬成 `Wireless` → **改到錯的商品**", CORAL),
    ("英文撇號炸掉 JS", "`'Didn't catch that'` ——英文化產生的撇號沒跳脫",
     "整份 JS 停擺：畫面正常但 WebSocket 從沒建立", CORAL),
]
ty2 = Inches(1.92); rh2 = Inches(0.94)
for i, (tag, what, effect, col) in enumerate(traps):
    y = ty2 + rh2 * i
    add_round(s, MX, y, Inches(11.87), Inches(0.84), fill=(LIGHT if i % 2 == 0 else WHITE),
              line=GREYE6, line_w=0.6)
    add_rect(s, MX, y, Inches(0.07), Inches(0.84), fill=col)
    add_text(s, MX + Inches(0.28), y + Inches(0.08), Inches(2.55), Inches(0.32),
             tag.replace("**", ""), size=12.5, bold=True, color=col)
    add_text(s, MX + Inches(0.28), y + Inches(0.44), Inches(2.55), Inches(0.32),
             f"坑 {i + 1}", size=10, color=GREY77)
    add_text(s, MX + Inches(3.0), y + Inches(0.08), Inches(4.4), Inches(0.34),
             what.replace("`", "").replace("**", ""), size=11, color=GREY44)
    add_text(s, MX + Inches(3.0), y + Inches(0.44), Inches(4.4), Inches(0.32),
             "↓", font=FONT_EN, size=9, color=GREYBB)
    add_text(s, MX + Inches(7.55), y + Inches(0.2), Inches(4.1), Inches(0.5),
             effect.replace("**", ""), size=11, bold=True,
             color=(CORAL if col == CORAL else GREY44), anchor=MSO_ANCHOR.MIDDLE)
add_round(s, MX, Inches(6.62), Inches(11.87), Inches(0.62), fill=DARK, shadow=True)
add_rich(s, MX + Inches(0.4), Inches(6.72), Inches(11.1), Inches(0.42),
         [[{"text": "找法  ", "size": 12.5, "bold": True, "color": TEAL},
           {"text": "逐句追 log 看實際執行路徑，不要只看輸入輸出猜",
            "size": 12, "bold": True, "color": WHITE},
           {"text": "——log 常見 clf 判對、模型也判對，卻被中文導向的守衛改掉。共歸納 19 類坑。",
            "size": 11.5, "color": GREYBB}]],
         anchor=MSO_ANCHOR.MIDDLE)
set_notes(s, "★這頁是給技術評審看的「移植的真實成本」。一般人以為做英文版＝翻譯 UI + 換模型，"
             "實際上**真正的工作量在散落各處的語言假設**——三個月為中文磨出來的規則層，每一條都"
             "藏著隱含假設。列五個最有代表性的（共歸納 19 類）：①21 個詞表全中文，英文句一個都不"
             "命中，被判成「只有商品名沒動作」全部轉 clarify。②更隱形的一類：對照表的**鍵**是中文"
             "（值反而是英文 slug），6 個類別 5 個查不到，而系統的導覽訊息還在教訪客這樣問。"
             "③門檻用中文字元數：英文字元數是中文 2-3 倍，'alert me when earphones drop below 30' "
             "才 7 個詞卻 31 字元，正常句被當長句。④最危險的一類——**演算法本身假設中文形態**："
             "用 split()[0] 剝規格尾巴，中文沒空白所以取第一段＝取全名，英文商品名全用空白分隔，"
             "`Wireless Mouse` 被腰斬成 `Wireless`，下一句追問就改到耳機的安全庫存＝**寫錯資料**。"
             "這種 bug 讀程式碼很難看出來（邏輯本身沒錯），只有跨句對話測試會暴露。⑤英文獨有的："
             "英文化產生的撇號沒跳脫，整份 JS 語法錯誤停擺，但 HTML/CSS 照常渲染——畫面看起來正常、"
             "實際 WebSocket 從沒建立過。中文版永遠不會有這個 bug（中文沒撇號）。"
             "方法論：**逐句追 log 看實際執行路徑**，不要從症狀猜成因。")
pn(s)


# ─── S14c 英文版 · 收斂成果（守衛 100%）★ ─────────────────────────
s = slide_blank()
title_bar(s, "ENGLISH BUILD · 收斂", f"從 651 到 {GUARD_EN}/{GUARD_EN}：把「可靠」再證明一次")
add_text(s, MX, Inches(1.35), Inches(11.8), Inches(0.4),
         f"英文版建了獨立的 {GUARD_EN} 句守衛庫（不是翻譯中文守衛——英文有自己的邊界：錯字型態、俗稱、閒聊搗蛋）。",
         size=13, color=GREY55)
kpi_row(s, Inches(1.9), [
    (f"{GUARD_EN}/{GUARD_EN}", "英文守衛通過率 100%"),
    ("19 類", "移植坑歸納"),
    ("25 個", "view 逐一看過畫面"),
    ("0", "已知未修破口"),
])
# 收斂曲線（分數演進）
add_text(s, MX, Inches(3.35), Inches(11.8), Inches(0.34),
         "收斂過程：每一次跳動都是一批結構性 bug 被找出來", size=13, bold=True, color=TEALDK)
steps = [("651", "首跑"), ("873", "13 輪規則層英文化"), ("891", "錯字長尾靠修復層"),
         (f"{GUARD_EN}", "詞典把關收尾")]
bw = Inches(2.72); bgap = Inches(0.32)
for i, (sc, lb) in enumerate(steps):
    x = MX + (bw + bgap) * i
    last = (i == len(steps) - 1)
    add_round(s, x, Inches(3.78), bw, Inches(1.12),
              fill=(TEALBG if last else LIGHT), line=(TEAL if last else GREYE6),
              line_w=(1.4 if last else 0.6), shadow=True)
    add_text(s, x, Inches(3.92), bw, Inches(0.5), sc, font=FONT_EN,
             size=26, bold=True, color=(TEALDK if last else GREY44),
             align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.12), Inches(4.46), bw - Inches(0.24), Inches(0.36),
             lb, size=10.5, color=GREY55, align=PP_ALIGN.CENTER)
    if not last:
        add_text(s, x + bw - Inches(0.04), Inches(4.16), Inches(0.34), Inches(0.4),
                 "›", font=FONT_EN, size=17, bold=True, color=GREYBB,
                 align=PP_ALIGN.CENTER)
# 底部：兩個方法論收穫
add_round(s, MX, Inches(5.15), Inches(5.82), Inches(1.42), fill=LIGHT, shadow=True)
add_text(s, MX + Inches(0.28), Inches(5.28), Inches(5.3), Inches(0.32),
         "「救不了」要有證據", size=13, bold=True, color=TEALDK)
add_text(s, MX + Inches(0.28), Inches(5.62), Inches(5.3), Inches(0.82),
         "兩次把「試過一種做法沒成功」寫成「這類問題無解」。"
         "最後一句錯字靠系統內建英文詞典區分「真詞 vs 錯字」修掉——"
         "訊號一直都在，只是沒接上。",
         size=11, color=GREY44, line_spacing=1.25)
add_round(s, MX + Inches(6.05), Inches(5.15), Inches(5.82), Inches(1.42), fill=LIGHT, shadow=True)
add_text(s, MX + Inches(6.33), Inches(5.28), Inches(5.3), Inches(0.32),
         "看畫面才算審完", size=13, bold=True, color=TEALDK)
add_text(s, MX + Inches(6.33), Inches(5.62), Inches(5.3), Inches(0.82),
         "守衛全綠不等於畫面對。改用截圖逐一看 25 個 view，"
         "抓到 JSON 完全看不到的破口：卡片承諾「說 delete AL001」，"
         "照打卻回查無此商品。",
         size=11, color=GREY44, line_spacing=1.25)
add_text(s, MX, Inches(6.76), Inches(11.87), Inches(0.4),
         "中英雙版並存：RPi5 同時跑 8002（英文，展場主力）與 8001（中文備援），開機自啟、訪客點分頁切換",
         size=12, bold=True, color=GREY55, align=PP_ALIGN.CENTER)
set_notes(s, "★英文版收斂成果頁。重點一：英文守衛庫是**重新建的 892 句**，不是把中文守衛翻譯過來"
             "——中文守衛大量測注音、同音字，英文沒有對應物；英文有自己的邊界（錯字型態、俗稱別名、"
             "英文閒聊搗蛋）。重點二：收斂曲線 651→873→891→892，每次跳動都是一批結構性 bug 被找出來，"
             "不是慢慢磨上去的。左下角這個教訓值得講：過程中**兩次**把『試過一種做法沒成功』寫成"
             "『這類問題無解』——第一次是 19 句雙錯字長尾，後來拆解發現是 8 個獨立的結構性 bug；"
             "第二次是最後那一句 `do we have scks`，記錄成『需要英文詞典依賴、救不了』，"
             "實際上樹莓派系統**內建**英文詞典，而且真正的成因根本不是字元相似度，是被守門員擋在"
             "門外。訊號一直都在，只是沒接上。右下角：守衛全綠≠畫面對，改用截圖逐一看過 25 個 view，"
             "抓到 JSON 看不到的破口——最典型的是卡片上明明寫著『To remove, say delete AL001』，"
             "訪客照打卻回『查無此商品』。最後一行：中英雙版在同一台樹莓派上並存，開機都自啟，"
             "訪客點瀏覽器分頁就能切語言。")
pn(s)


# ─── S13a 語音 POC · 全鏈架構 ────────────────────────────────
s = slide_blank()
title_bar(s, "VOICE POC · 全離線語音輸入", "訪客用「講的」查倉管，ASR 全程跑在樹莓派")
add_text(s, MX, Inches(1.42), Inches(11.8), Inches(0.4),
         "打字對展場訪客不夠自然——加一層離線語音輸入，同一套倉管後端不動，語音只是新入口。",
         size=13.5, color=GREY55)
# 縱向全鏈流程
vchain = [
    ("前端錄音", "Siri 式：點一下 → 講 → 靜音自動結束（瀏覽器 VAD 偵測）", TEAL, "🎙️"),
    ("whisper.cpp", f"OpenAI whisper，RPi5 CPU 純離線辨識（中英同顆 {ASR['name']}，{ASR['lat']}/句）", TEALDK, "🧠"),
    ("OpenCC 轉繁", "簡體 → 繁體，順便轉台灣用語（僅中文版需要）", NAVY, "🔄"),
    ("同音修正層", "倉別 / 動詞 / 量詞 / 異體字——只掛 ASR 出口，不碰倉管核心", AMBER, "🔧"),
    ("倉管 WS", "既有守衛庫 + 發音容錯層接手，回答與打字完全一致", TEALDK, "📦"),
]
vy = Inches(2.05)
for i, (t, d, col, ic) in enumerate(vchain):
    y = vy + Inches(0.94) * i
    add_round(s, MX, y, Inches(11.87), Inches(0.76), fill=LIGHT, shadow=True)
    add_icon_circle(s, MX + Inches(0.28), y + Inches(0.13), Inches(0.5), ic,
                    circle=col, gcolor=WHITE, gsize=15)
    add_text(s, MX + Inches(1.05), y + Inches(0.06), Inches(3.0), Inches(0.64),
             t, size=15, bold=True, color=col, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, MX + Inches(4.05), y + Inches(0.06), Inches(7.6), Inches(0.64),
             d, size=12.5, color=GREY44, anchor=MSO_ANCHOR.MIDDLE)
    if i < len(vchain) - 1:
        add_arrow(s, Inches(6.5), y + Inches(0.78), Inches(0.36), Inches(0.14), fill=GREYBB)
set_notes(s, "語音 POC 全鏈架構（2026-07 新增）。核心設計：語音只是新入口，倉管後端"
             "完全不動。前端 Siri 式點一下自動結束；whisper.cpp 在 RPi5 純 CPU 離線"
             "辨識（展場無網路也能跑）；OpenCC 轉繁（僅中文版）+ 出口修正層清理 ASR 錯字；最後交回"
             "既有倉管 WS。同音修正層刻意只掛 /api/asr 出口——打字訪客零影響、守衛零風險。")
pn(s)


# ─── S13a2 語音 POC · ASR 選型（全面改用歐美模型）★ ────────────────
s = slide_blank()
title_bar(s, "VOICE POC · 選型", "壓縮程度怎麼選：實測打敗直覺")
add_text(s, MX, Inches(1.38), Inches(11.8), Inches(0.4),
         "選型鐵律：① 供應鏈來源可控（歐美模型）② RPi5 CPU 純離線跑 ③ 中英共用同一套 whisper.cpp runtime。"
         "同一顆 small 模型有三種壓縮程度可選，實測結果推翻了兩個直覺。",
         size=13, color=GREY55)
# 對照表：欄位 = 模型 / 體積 / 語言 / RPi5 延遲 / 準確度 / 判定
col_x = [MX, Inches(3.05), Inches(4.35), Inches(6.05), Inches(7.95), Inches(10.05)]
col_w = [Inches(2.25), Inches(1.45), Inches(1.85), Inches(2.05), Inches(2.25), Inches(2.55)]
headers = ["模型", "體積", "語言", "RPi5 延遲", "純模型辨識率", "本專案判定"]
ty = Inches(1.95); th = Inches(0.5)
# 表頭
add_round(s, MX, ty, Inches(11.87), th, fill=DARK)
for i, hd in enumerate(headers):
    add_text(s, col_x[i] + Inches(0.12), ty + Inches(0.09), col_w[i] - Inches(0.2), Inches(0.36),
             hd, size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
# 2026-08-06 全面重量：所有數字改用**當日同一套方法**實測（舊值作廢）。
#   延遲＝英文版·機二·走網頁真實路徑·只計 ASR 段（不含 270M 推論）。
#   準確度＝端到端含容錯層（訪客實際體驗），100 句 × 三層噪音平均。
# ⚠️ 準確度欄一律是**純模型辨識率**（無容錯層）——選型比的是模型本身；
#   混入容錯層會變成比「模型＋工程」，模糊掉判斷依據。
#   容錯層接手後的成品表現見「100 句最終實測」頁。
rows = [
    ("whisper small + 量化 q8", "252 MB", "多語（中英同檔）", "2.5s",
     "真人 39% / 合成 87%", "★ 定案：準度夠又壓得下來", True),
    ("whisper small 未量化", "465 MB", "多語（中英同檔）", "10 秒以上",
     "同上（量化不損準度）", "準度夠但太慢 ⇒ 量化解決", False),
    ("whisper base", "141 MB", "多語（含中文）", "2.2s",
     "真人明顯較差", "× 準度不足，升級 small", False),
    ("whisper tiny.en", "74 MB", "英文專用", "0.9s",
     "合成音好、真人差", "× 台灣腔真人辨識不堪用", False),
]
ry = ty + th
# 2026-08-06：加了 q8_0 與全精度兩列（5→6 列）。原列高 0.68 會讓表格底部
#   落到 y=6.53，壓到下方 y=5.98 的深色摘要框 ⇒ 列高收到 0.55，
#   表格底部回到 6.00 附近，與摘要框不重疊。
rh = Inches(0.55)
for r, (m, p, rt, off, cer, verd, chosen) in enumerate(rows):
    y = ry + rh * r
    bg = TEALBG if chosen else (LIGHT if r % 2 else WHITE)
    add_rect(s, MX, y, Inches(11.87), rh, fill=bg, line=GREYE6, line_w=0.5)
    add_text(s, col_x[0] + Inches(0.12), y + Inches(0.04), col_w[0] - Inches(0.2), Inches(0.47),
             m, size=12.5, bold=chosen, color=(TEALDK if chosen else DARK), anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, col_x[1] + Inches(0.12), y + Inches(0.04), col_w[1] - Inches(0.2), Inches(0.47),
             p, font=FONT_EN, size=12, bold=chosen, color=GREY44, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, col_x[2] + Inches(0.12), y + Inches(0.04), col_w[2] - Inches(0.2), Inches(0.47),
             rt, size=11, color=GREY44, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, col_x[3] + Inches(0.12), y + Inches(0.04), col_w[3] - Inches(0.2), Inches(0.47),
             off, font=FONT_EN, size=11.5, color=(TEALDK if chosen else GREY55), bold=chosen,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, col_x[4] + Inches(0.12), y + Inches(0.04), col_w[4] - Inches(0.2), Inches(0.47),
             cer, size=11.5, color=GREY44, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, col_x[5] + Inches(0.12), y + Inches(0.04), col_w[5] - Inches(0.2), Inches(0.47),
             verd, size=11.5, bold=chosen, color=(TEALDK if chosen else GREY55), anchor=MSO_ANCHOR.MIDDLE)
# 底部：兩段結論（2026-08-06 user 定調敘事——不提其他量化版本，簡潔有力）
add_round(s, MX, Inches(5.98), Inches(11.87), Inches(1.18), fill=DARK, shadow=True)
add_rich(s, MX + Inches(0.4), Inches(6.12), Inches(11.1), Inches(0.5),
         [[{"text": "選型邏輯  ", "size": 13, "bold": True, "color": TEAL},
           {"text": "small 夠準但太慢 → 量化壓縮 → q8_0", "size": 13, "bold": True, "color": WHITE},
           {"text": "——小模型（tiny/base）準度不足；small 準度夠卻要 10 秒以上，"
                    "展場站著等不了。量化把它壓到 252 MB、2.5 秒，準度幾乎不損。",
            "size": 12, "color": GREYBB}]],
         anchor=MSO_ANCHOR.MIDDLE)
add_rich(s, MX + Inches(0.4), Inches(6.63), Inches(11.1), Inches(0.42),
         [[{"text": "為何純模型 39% 仍可用  ", "size": 13, "bold": True, "color": AMBER},
           {"text": "判準是「一字不差」才算對", "size": 12.5, "bold": True, "color": WHITE},
           {"text": "——實際字元正確率 89%（每 10 字才錯 1 字），多為 earphones→earphone "
                    "這類小差異。容錯層接手後端到端 80%，那才是訪客體驗到的。",
            "size": 12, "color": GREYBB}]],
         anchor=MSO_ANCHOR.MIDDLE)
set_notes(s, "【2026-08-06 定案：small + 量化 q8_0，中英同一顆】\n"
             "① 為何是 small：tiny.en 與 base 對**台灣腔真人**明顯不夠——tiny 在合成音"
             "表現不錯，換真人就掉下來（選型不能只信 TTS）。small 才撐得住。\n"
             "② 為何要量化：small 未量化 465MB、每句 10 秒以上，展場站著等不了。"
             "量化壓到 252MB、2.5 秒，準度幾乎不損 ⇒ 用 q8_0。\n"
             "③ 為何是多語版不是英文專用版：多語模型訓練時看過大量非母語者講的英文，"
             "對台灣腔更寬容；英文專用版聽慣標準英美腔反而不適應。中英共用同一顆檔案，"
             "只靠 -l 旗標切語言，不必維護兩套框架。\n"
             "④ 純模型 39% 看起來低，是因為判準「一字不差」。字元正確率 89%，"
             "容錯層接手後端到端 80% —— 那才是訪客體驗到的數字。\n\n"
             "★語音選型頁（技術評審向）。**這一頁 2026-07-27 全面改版**：原本用阿里的 "
             "Fun-ASR-Nano，後來定調**只用歐美模型**（供應鏈來源可控），中英兩版都換成 "
             "OpenAI whisper.cpp。換完的額外好處：中英共用同一套 runtime，不必維護兩套框架。"
             "選型結果：英文版 tiny.en（74MB / 0.94s / WER 9.3%）、中文版 base（141MB / 2.15s）。"
             "兩個值得講的洞見——①**反直覺**：base.en 比 tiny.en 更大更慢，WER 反而略差（10.2% vs "
             "9.3%），因為倉管查詢句短、句型固定，tiny 的容量已經夠用，變大的收益顯現不出來；"
             "這推翻『模型越大越好』的直覺，也呼應整個專案『選對尺寸勝過選大尺寸』的主張。"
             "②**選型不能只看 WER**：中文版 base 真人 8 句字面只對 3 句，但端到端 8/8 全部答對——"
             "聽錯的『盡量啤酒酷醇』『無限滑鼠擴存』『瑜伽店』全被文字端容錯層救回。這正好證明"
             "前面幾頁講的容錯層是真的在扛事情。代價誠實講：中文 base 2.15s/句比原本略慢，但中文版"
             "是備案、英文版才是展場主力（tiny.en 0.94s 很順）。英文數據為 TTS 合成音實測"
             "（5 腔調×20 句×3 噪音層），中文為 user 錄的真人音——TTS 是下限估計，評審追問據實說明。\n"
             "**2026-08 演進（本頁現役列）**：英文先升級 small-q5_0＋-ac 640（比 tiny.en 更準且 "
             "3.45s/句——比 tiny.en 慢但台灣腔通過率 27%→60%，速度靠 -ac 640 從 6.7s 拉回）；8/5 手機麥克風實測"
             "暴露 base 中文極限（『藍牙耳機庫存』聽成音節碎片、容錯層救不動），以 user 自錄真人 "
             "100 句重跑基準：base 引擎 2.10s/CER中位 0.33、small 3.54s/0.25；**公平對稱端到端"
             "（同無修正層）base 36/100 vs small 47/100；small 掛回容錯層 66/100**——慢 1.4 秒買 "
             "+30 分答對率，且容錯層在 small 的詞形輸出上多救 19 分（47→66），驗證『輸出要像詞、"
             "救援鏈才接得上手』。中英自此統一同一顆 175MB 模型檔，只差 -l 語言旗標。")
pn(s)


# ─── S13a3 語音 POC · 部署架構（單一模型 · 中英雙軌）★ ────────────
s = slide_blank()
title_bar(s, "VOICE POC · 部署架構", "從三個模型變成一個：換型順帶把架構砍薄")
add_text(s, MX, Inches(1.35), Inches(11.8), Inches(0.36),
         "原方案是 VAD＋Encoder＋LLM 解碼器三顆模型（910 MB）。換 whisper 後：端點偵測交給瀏覽器、編碼解碼合成一顆模型檔。",
         size=12.5, color=GREY55)
# 三張卡：runtime + 英文模型 + 中文模型
comp = [
    ("端點偵測移到前端", "Web Audio API · 零模型", "RMS 音量偵測",
     "瀏覽器即時算音量：靜音 1.2s 自動送出、15s 硬上限；省掉一顆 VAD 模型", NAVY, "偵"),
    ("編碼＋解碼合一", ASR['file'], f"{ASR['size']} · 中英共用一顆",
     "whisper 端到端 encoder-decoder：聲學編碼與文字解碼同一檔；-ac 640 削掉短句用不到的空白運算", TEALDK, "聲"),
    ("同一套 runtime", "-l en ／ -l zh", "語言只是一個參數",
     "中英只差語言旗標、共用同一顆模型檔；多語權重輸出簡體，中文仍走 OpenCC 轉繁", TEAL, "文"),
]
cw = Inches(3.83); ch = Inches(2.55); cgap = Inches(0.19)
cy = Inches(2.0)
for i, (name, fn, size, desc, col, gl) in enumerate(comp):
    x = MX + (cw + cgap) * i
    add_round(s, x, cy, cw, ch, fill=LIGHT, shadow=True)
    add_icon_circle(s, x + Inches(0.28), cy + Inches(0.24), Inches(0.62), gl,
                    circle=col, gcolor=WHITE, gsize=18)
    add_text(s, x + Inches(1.05), cy + Inches(0.26), cw - Inches(1.2), Inches(0.4),
             name, size=13.5, bold=True, color=col, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(1.05), cy + Inches(0.62), cw - Inches(1.2), Inches(0.3),
             size, font=FONT_EN, size=10.5, color=GREY77, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.3), cy + Inches(1.08), cw - Inches(0.55), Inches(0.34),
             fn, font=FONT_EN, size=10.5, bold=True, color=DARK)
    add_text(s, x + Inches(0.3), cy + Inches(1.46), cw - Inches(0.55), Inches(1.0),
             desc, size=11.5, color=GREY44, line_spacing=1.2)
    if i < 2:
        add_text(s, x + cw - Inches(0.02), cy + ch / 2 - Inches(0.18), Inches(0.24),
                 Inches(0.36), "▶", font=FONT_EN, size=13, color=GREYBB,
                 align=PP_ALIGN.CENTER)
# 部署管線（橫向一條 pipeline）
py = Inches(4.85)
add_text(s, MX, py - Inches(0.06), Inches(11.8), Inches(0.34),
         "實際部署管線（POST /api/asr，全程本機、無網路）", size=13, bold=True, color=TEALDK)
pipe = ["前端錄音 + VAD\n靜音 1.2s 自動停", "ffmpeg\n轉 16k mono",
        f"whisper-cli\n{ASR['name']} · -ac 640",
        "OpenCC\n簡→繁（僅中文）", "出口正規化\n大小寫/錯字", "倉管 WS\n回答"]
pw = Inches(1.78); ph = Inches(0.82); pgap = Inches(0.19)
px0 = MX
for i, step in enumerate(pipe):
    x = px0 + (pw + pgap) * i
    isc = (i == 2)  # 三元件那格強調
    add_round(s, x, py + Inches(0.36), pw, ph, fill=(TEALBG if isc else LIGHT),
              line=(TEAL if isc else None), line_w=1.2, shadow=True)
    add_text(s, x + Inches(0.08), py + Inches(0.44), pw - Inches(0.16), ph - Inches(0.16),
             step, size=10.5, bold=isc, color=(TEALDK if isc else GREY44),
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    if i < len(pipe) - 1:
        add_text(s, x + pw - Inches(0.02), py + Inches(0.5), Inches(0.22), Inches(0.5),
                 "›", font=FONT_EN, size=16, bold=True, color=GREYBB, align=PP_ALIGN.CENTER)
# 底部部署事實條（2026-08-06 加入 -ac 640 的時間管控——user 定調：要讓老闆
#   知道速度是設計出來的，且語音長度有上限、不是無限吃）
add_round(s, MX, Inches(6.28), Inches(11.87), Inches(0.95), fill=DARK, shadow=True)
add_rich(s, MX + Inches(0.35), Inches(6.38), Inches(11.2), Inches(0.36),
         [[{"text": "RPi5 部署  ", "size": 12.5, "bold": True, "color": TEAL},
           {"text": f"910 MB 三顆 → {ASR['size']} 一顆（中英共用）｜VAD 零模型｜"
                    f"英文 {ASR['lat']}/句｜純 CPU、零 GPU、零雲端",
            "size": 11.5, "color": GREYBB}]],
         anchor=MSO_ANCHOR.MIDDLE)
add_rich(s, MX + Inches(0.35), Inches(6.78), Inches(11.2), Inches(0.36),
         [[{"text": "時間管控  ", "size": 12.5, "bold": True, "color": AMBER},
           {"text": "-ac 640 = 一次只吃 12.8 秒（預設 30 秒）",
            "size": 11.5, "bold": True, "color": WHITE},
           {"text": "——問句 2-3 秒，其餘都是空白運算；削掉後速度快一倍，"
                    "同時等於替語音輸入設了長度上限。",
            "size": 11.5, "color": GREYBB}]],
         anchor=MSO_ANCHOR.MIDDLE)
set_notes(s, "★語音部署架構頁（回應「語音這塊怎麼部署」「VAD 和 Encoder 現在怎麼做」）。"
             "**2026-07-27 隨 ASR 換型改版**：原本是 Fun-ASR 的三顆模型串起來跑——VAD（fsmn-vad "
             "1.6MB，判斷何時在講話）＋Encoder（447MB，聲學編碼）＋Qwen3-0.6B 解碼器（461MB，"
             "生成文字），合計 910MB。換成 whisper 之後這三塊各自的去向如下，這是評審會追問的重點：\n"
             "① **VAD → 移到前端瀏覽器，不用模型了**。whisper.cpp 本身有 VAD 支援（--vad 搭配 "
             "Silero 模型），但我們**沒有啟用**；改用 Web Audio API 在瀏覽器即時算 RMS 音量："
             "偵測到開始講話後，靜音持續 1.2 秒就自動送出，另有 15 秒硬上限（展場吵雜時 VAD 可能"
             "永遠不收尾）。好處是零模型負擔、零延遲、麥克風端就地判斷；也是訪客體感「Siri 式"
             "點一下自動結束」的來源。\n"
             "③ **-ac 640 是速度與長度上限的關鍵**（2026-08-06 補充）：whisper 預設一次處理 "
             "30 秒音訊（audio-ctx 1500），不足會補靜音湊滿——倉管問句只有 2-3 秒，等於九成"
             "算力花在算空白。削到 640（12.8 秒）速度快一倍且準度不掉。**副作用正好是優點**："
             "它同時替語音輸入設了 12.8 秒上限，訪客講再長也不會無限吃資源。實測錄音最長 "
             "3.34 秒，12.8 秒留了近四倍餘裕；曾評估再削到 6.4 秒，但長句停頓會被截斷，"
             "定調維持 640。中英兩版都是這個值（英文走環境變數、中文寫在命令列）。\n"
             "② **Encoder 與解碼器 → 合併成同一個模型檔**。whisper 是端到端的 encoder-decoder "
             "架構，聲學編碼與文字解碼本來就在同一份權重裡，不像 Fun-ASR 要三個檔案分開載入。"
             "所以現在只有一顆 .bin，中英兩版只差 -m 參數指到 tiny.en 還是 base。\n"
             "管線：前端錄 webm（含 VAD 自動停）→ ffmpeg 轉 16k mono → whisper-cli → OpenCC 轉繁"
             "（**只有中文版需要**，英文版已移除）→ 出口正規化（英文做大小寫攤平、中文做同音錯字"
             "修正）→ 交倉管 WS，全程本機無網路。**模型體積 910MB → 215MB、英文延遲 2.5s → 0.94s**"
             "——換掉來源不符的模型不但沒犧牲，反而更輕更快、架構更薄。\n"
             "**2026-08 更新**：中英已統一為同一顆 ggml-small-q5_0.bin（175MB 多語、量化 q5_0），"
             "兩版只差 -l en / -l zh 語言旗標——連『兩顆模型』都省成一顆，部署與還原再薄一層。"
             "速度靠 -ac 640（audio-ctx 從 30 秒上下文削到短句實際需要）守住：EN 3.45s、ZH 3.54s/句（tiny.en 時代的 0.94-1.1s 是舊引擎數字，換 small 後中英同級）。")
pn(s)


# ─── S13b 語音 POC · 三環境噪音測試 ─────────────────────────
s = slide_blank()
title_bar(s, "VOICE POC · 展場噪音實測", "念一次真人聲，自動測三種環境")
add_text(s, MX, Inches(1.42), Inches(11.8), Inches(0.4),
         "方法突破：真人念一次乾淨版並存錄音，用同一份錄音自動混入賣場人潮噪音重測——對比最公平。",
         size=13.5, color=GREY55)
# KPI：三環境通過率
kpi_row(s, Inches(2.15), [
    ("78%", "乾淨（正常音量）"),
    ("77%", "一般人潮"),
    ("73%", "尖峰吵雜"),
], box_h=1.3, num_size=40)
# 關鍵發現卡片
finds = [
    ("噪音幾乎零影響", "一般展場 −1%、尖峰只 −5%；純被噪音壓垮的僅 5 句", TEAL),
    ("音量才是關鍵", "小聲時摩擦音（ㄕ/ㄘ）糊掉→亂猜；正常音量直接解決大半", AMBER),
    ("失敗多為 ASR 極限", "整詞聽錯，訪客看辨識文字重講即可，非系統 bug", NAVY),
]
fy = Inches(3.85)
for i, (t, d, col) in enumerate(finds):
    y = fy + Inches(0.95) * i
    add_round(s, MX, y, Inches(11.87), Inches(0.8), fill=LIGHT, shadow=True)
    add_circle(s, MX + Inches(0.32), y + Inches(0.28), Inches(0.24), col)
    add_text(s, MX + Inches(0.82), y + Inches(0.08), Inches(3.8), Inches(0.66),
             t, size=15, bold=True, color=col, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, MX + Inches(4.7), y + Inches(0.08), Inches(6.9), Inches(0.66),
             d, size=12.5, color=GREY44, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, MX, Inches(6.75), Inches(11.8), Inches(0.5),
         "結論：webcam + 270M + 正常音量＝展場夠用，不需升級模型或麥克風。",
         size=14, bold=True, color=TEALDK)
set_notes(s, "★語音展場可行性的量化證據。三環境通過率 78/77/73%——關鍵訊息是"
             "「噪音幾乎不影響」（light −1%、heavy −5%）。方法上的巧思：真人只念一次、"
             "存下錄音，之後自動混噪重測，念一次測三種環境、對比公平。三大發現：噪音不是"
             "問題、音量才是關鍵變數、失敗幾乎全是 ASR 模型極限（訪客重講可解）。"
             "最終結論：現有硬體（webcam + 270M）＋正常音量就夠展場用，不用花錢升級。")
pn(s)


# ─── S13c 語音 POC · 兩大容錯層 ─────────────────────────────
s = slide_blank()
title_bar(s, "VOICE POC · 容錯設計", "兩層容錯：拼音智慧比對 + 固定規則修正")
add_text(s, MX, Inches(1.42), Inches(11.8), Inches(0.4),
         "270M 辨識不完美（滑鼠→華數/華族）。主角是「拼音容錯層」——把錯字轉拼音智慧比對；再配一層固定規則收尾。",
         size=13.5, color=GREY55)
# 左卡：發音容錯層（拼音，技術主角 → 較寬）
lx = MX
add_round(s, lx, Inches(2.15), Inches(6.35), Inches(3.75), fill=TEALBG, shadow=True)
add_rich(s, lx + Inches(0.35), Inches(2.38), Inches(5.65), Inches(0.42),
         [[{"text": "① 拼音容錯層", "size": 17, "bold": True, "color": TEALDK},
           {"text": "  · 主角 · 智慧比對", "size": 12, "bold": True, "color": TEAL}]])
add_text(s, lx + Inches(0.35), Inches(2.85), Inches(5.65), Inches(0.4),
         "字形救不到 → 轉拼音、滑窗比對商品名", font=FONT_ZH, size=12.5, bold=True, color=DARK)
pfx = [
    "同音字形遠：滑鼠 vs 華數，字形 0 分",
    "轉拼音一比即中：huashu ≈ huashu",
    "音節還原 zu→zhu／su→shu，救捲舌混淆",
    "滑窗掃句：核心名拼音是句拼音子串 → 命中",
    "門檻 0.82 防誤配（衛生棉≠衛生紙）",
    "字形優先、拼音救底 → 打字＋守衛零回歸",
]
for i, t in enumerate(pfx):
    y = Inches(3.32) + Inches(0.42) * i
    add_text(s, lx + Inches(0.35), y, Inches(0.3), Inches(0.4), "·", size=15,
             bold=True, color=TEAL)
    add_text(s, lx + Inches(0.62), y, Inches(5.5), Inches(0.4), t, size=12,
             color=GREY44)
# 右卡：同音修正層（固定規則 → 較窄）
rx2 = Inches(7.55)
rw2 = SLIDE_W - rx2 - MX
add_round(s, rx2, Inches(2.15), rw2, Inches(3.75), fill=LIGHT, shadow=True)
add_text(s, rx2 + Inches(0.32), Inches(2.38), rw2 - Inches(0.6), Inches(0.42),
         "② 固定規則修正", size=17, bold=True, color=AMBER)
add_text(s, rx2 + Inches(0.32), Inches(2.85), rw2 - Inches(0.6), Inches(0.4),
         "掛 ASR 出口，不碰倉管核心", font=FONT_ZH, size=12.5, bold=True, color=DARK)
sfx = [
    "倉別：總/藏/昌倉 → 中/北倉",
    "動詞：近→進、谷→補",
    "量詞：臺→台（OpenCC 差異）",
    "異體字：溼→濕、賬→帳、周→週",
    "退貨中文數字豁免",
    "打字訪客零影響、守衛零風險",
]
for i, t in enumerate(sfx):
    y = Inches(3.32) + Inches(0.42) * i
    add_text(s, rx2 + Inches(0.32), y, Inches(0.3), Inches(0.4), "·", size=15,
             bold=True, color=AMBER)
    add_text(s, rx2 + Inches(0.58), y, Inches(4.0), Inches(0.4), t, size=12,
             color=GREY44)
add_text(s, MX, Inches(6.3), Inches(11.8), Inches(0.5),
         "分工：拼音層救「沒見過的音近錯」（會智慧比對）；規則層收「反覆出現的固定錯」（查表最快最穩）。下一頁看拼音層怎麼一步步救。",
         size=12.5, bold=True, color=TEALDK)
set_notes(s, "語音容錯兩層設計，都是真人聲實測一句句磨出來的。刻意分工：左＝拼音容錯層（技術主角）"
             "——處理「沒見過的音近錯」，靠智慧比對而非查表：ASR 錯字多是「同音但字形差很遠」"
             "（滑鼠→華數），字形比對救不到、轉拼音就中；做了捲舌平舌音節還原（zu→zhu、su→shu）；"
             "用滑窗掃整句找商品核心名拼音；門檻 0.82（太低會把衛生棉誤配衛生紙）。右＝固定規則修正"
             "——處理「反覆出現的固定錯」（倉別/動詞/量詞/異體字），這種錯每次都一樣，直接查表最快最穩，"
             "不需要智慧比對。兩層都只掛 ASR 出口、不碰 warehouse.py，所以打字訪客零影響、守衛零風險。"
             "下一頁用流程圖把拼音層的五步決策拆開講。")
pn(s)


# ─── S13c2 語音 POC · 拼音修正流程圖 ★老闆愛看圖 ──────────────────
s = slide_blank()
title_bar(s, "VOICE POC · 拼音修正流程", "一個錯字怎麼被救回：字形先試，拼音救底")
add_text(s, MX, Inches(1.35), Inches(11.8), Inches(0.36),
         "核心：ASR 錯字多是「音同、字形零重疊」——字形比對必敗，轉成拼音一比即中。全程純字串運算，RPi5 零負擔。",
         size=12.5, color=GREY55)
# ── 左：直向決策流程（5 步 + 判斷菱形感）──
fx = MX
fw = Inches(6.7)
steps = [
    ("輸入 keyword", "270M 抽出的商品詞（常髒／殘／錯）", NAVY, "in"),
    ("① 字形比對優先", "先走 LCS 字形比對，命中就用 → 不進拼音層", TEALDK, "step"),
    ("② 三道排除閘", "含寫入動詞／倉別／數字 → 不救；>6 字 → 不救", AMBER, "gate"),
    ("③ 轉拼音 + 音節還原", "lazy_pinyin，捲舌平舌還原（zu→zhu、su→shu）", TEAL, "step"),
    ("④ 滑窗比對商品拼音", "核心名拼音 vs 句拼音，difflib 取最佳對齊", TEAL, "step"),
    ("⑤ 門檻 0.82 判定", "≥0.82 救回；否則回空 → 交 clarify 反問", TEALDK, "dec"),
]
sy = Inches(1.95); sh = Inches(0.66); sgap = Inches(0.185)
for i, (name, desc, col, kind) in enumerate(steps):
    y = sy + (sh + sgap) * i
    shp = add_round(s, fx, y, fw, sh, fill=(LIGHT if kind != "in" else col), shadow=True)
    tcol = WHITE if kind == "in" else col
    add_text(s, fx + Inches(0.28), y + Inches(0.07), Inches(2.75), Inches(0.54),
             name, size=13.5, bold=True, color=tcol, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, fx + Inches(3.05), y + Inches(0.07), fw - Inches(3.3), Inches(0.54),
             desc, size=11, color=(GREYBB if kind == "in" else GREY44),
             anchor=MSO_ANCHOR.MIDDLE)
    if i < len(steps) - 1:
        add_text(s, fx + fw / 2 - Inches(0.2), y + sh - Inches(0.03), Inches(0.4),
                 Inches(0.22), "▼", font=FONT_EN, size=11, color=GREYBB,
                 align=PP_ALIGN.CENTER)
# ── 右：具體範例貫穿（同一句一路走下來）──
rx = MX + Inches(7.05)
rw = SLIDE_W - rx - MX
add_round(s, rx, Inches(1.95), rw, Inches(5.15), fill=DARK, shadow=True)
add_text(s, rx + Inches(0.32), Inches(2.15), rw - Inches(0.64), Inches(0.36),
         "實例：「北倉的滑鼠有多少」", size=14, bold=True, color=TEAL)
trace = [
    ("ASR 聽成", "北倉的華數有多少", CORAL),
    ("抽 keyword", "華數", GREYBB),
    ("① 字形比對", "華數 vs 滑鼠 → 0 重疊，失敗", GREYBB),
    ("② 排除閘", "無寫入詞／無數字／2 字 → 放行", GREYBB),
    ("③ 轉拼音", "華數 = huashu（su→shu 還原）", WHITE),
    ("④ 滑窗比對", "滑鼠 huashu ⊂ 句拼音 → 命中", WHITE),
    ("⑤ 門檻", "0.95 ≥ 0.82 → 救回", TEAL),
    ("結果", "→ 無線滑鼠 ✓ 答對", TEAL),
]
for i, (k, v, col) in enumerate(trace):
    y = Inches(2.62) + Inches(0.55) * i
    add_text(s, rx + Inches(0.32), y, Inches(1.5), Inches(0.5),
             k, size=11.5, bold=True, color=GREYBB, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, rx + Inches(1.9), y, rw - Inches(2.2), Inches(0.5),
             v, size=12, bold=(col in (TEAL, WHITE)), color=col, anchor=MSO_ANCHOR.MIDDLE)
set_notes(s, "★拼音修正流程圖（老闆/評審最有感的一頁）。左邊是決策流程，右邊用同一句"
             "「北倉的滑鼠有多少」一路走到底。核心洞見：ASR 錯字的特徵是「音同、字形零重疊」"
             "（華數 vs 滑鼠一個字都不重疊），所以字形比對一定失敗、轉成拼音一比即中。"
             "五步：①字形優先（命中就不進拼音層，零回歸的保證）②三道排除閘（含寫入動詞/倉別/"
             "數字的句子不救——這是當初打壞 63 條守衛後加的，只救乾淨短查詢詞）③轉拼音時做捲舌"
             "平舌音節還原（zu→zhu、su→shu，因為台灣國語與 ASR 常把滑鼠 huashu 聽成華族 huazu）"
             "④滑窗比對商品核心名拼音⑤門檻 0.82。門檻 0.82 是實測血淚調出來的——放寬到 0.78 會把"
             "「衛生棉」誤配「衛生紙」（同 0.80），所以不靠放門檻、靠精準音節還原來拉開分數。"
             "全程純字串運算，RPi5 零算力負擔。救不到就回空、交給系統反問，絕不亂猜。")
pn(s)


# ─── S13d 語音 POC · 聽錯→救回 實測範例 ★ ──────────────────────
s = slide_blank()
title_bar(s, "VOICE POC · 救回實例", "ASR 聽錯，容錯層照樣答對（真人聲實測）")
add_text(s, MX, Inches(1.38), Inches(11.8), Inches(0.4),
         "以下全是 2026-07 真人聲實測（webcam）的原始紀錄：ASR 明明聽錯，經容錯層修正後結果正確。",
         size=13, color=GREY55)
# 三欄對照表：原句 / ASR 聽成 / 救回結果
cx = [MX, Inches(4.55), Inches(8.35)]
cw = [Inches(3.7), Inches(3.6), Inches(4.05)]
hd = ["訪客原句", "ASR 聽成（錯）", "容錯層救回 → 答對"]
ty = Inches(1.95); th = Inches(0.46)
add_round(s, MX, ty, Inches(11.87), th, fill=DARK)
for i in range(3):
    add_text(s, cx[i] + Inches(0.14), ty + Inches(0.07), cw[i] - Inches(0.24), Inches(0.34),
             hd[i], size=12, bold=True, color=(TEAL if i == 2 else WHITE), anchor=MSO_ANCHOR.MIDDLE)
save_rows = [
    ("北倉進五十個滑鼠", "北藏近五十個華族", "藏→倉·近→進·華族→滑鼠"),
    ("幫我在北倉加五十個滑鼠", "…加五十個花束", "花束→滑鼠（拼音同音）"),
    ("藍牙耳機庫存", "藍芽耳機庫存", "藍芽→藍牙（OpenCC 用語）"),
    ("中倉衛生紙還有嗎", "中餐衛生紙還有嗎", "中餐→中倉"),
    ("運動壓縮臂套庫存", "運動壓縮筆套庫存", "筆套→臂套"),
    ("精釀啤酒庫存", "儘量啤酒庫存", "儘量→精釀"),
    ("防蚊液庫存", "防蚊衣庫存", "防蚊衣→防蚊液"),
    ("衛生紙的帳對不上", "衛生紙的賬對不上", "賬→帳（異體字）"),
    ("南倉出十五個瑜珈墊", "南倉出十五個瑜伽墊", "瑜伽→瑜珈"),
]
ry = ty + th; rh = Inches(0.435)
for r, (o, wrong, fix) in enumerate(save_rows):
    y = ry + rh * r
    add_rect(s, MX, y, Inches(11.87), rh, fill=(LIGHT if r % 2 else WHITE), line=GREYE6, line_w=0.5)
    add_text(s, cx[0] + Inches(0.14), y + Inches(0.05), cw[0] - Inches(0.24), Inches(0.34),
             o, size=12, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, cx[1] + Inches(0.14), y + Inches(0.05), cw[1] - Inches(0.24), Inches(0.34),
             wrong, size=12, color=CORAL, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, cx[2] + Inches(0.14), y + Inches(0.05), cw[2] - Inches(0.24), Inches(0.34),
             fix, size=12, bold=True, color=TEALDK, anchor=MSO_ANCHOR.MIDDLE)
# 底部拼音證據
add_round(s, MX, Inches(6.5), Inches(11.87), Inches(0.72), fill=TEALBG)
add_rich(s, MX + Inches(0.35), Inches(6.63), Inches(11.2), Inches(0.48),
         [[{"text": "為何拼音救得到：  ", "size": 12.5, "bold": True, "color": TEALDK},
           {"text": "華數 huashu = 滑鼠 huashu（完全同音）｜藍雅爾基 lanyaerji = 藍牙耳機 lanyaerji｜到齊 daoqi = 到期 daoqi",
            "font": FONT_EN, "size": 12, "color": GREY44},
           {"text": "  → 字形零重疊、拼音一比即中", "size": 12, "color": GREY44}]],
         anchor=MSO_ANCHOR.MIDDLE)
set_notes(s, "★聽錯→救回實例頁（技術評審最有感）。強調這些全是真人聲實測原始紀錄、不是虛構："
             "中欄紅字是 ASR 真的聽錯的字，右欄綠字是容錯層怎麼救回。最漂亮的一句「北倉進五十個滑鼠」"
             "一次踩三種錯（倉別藏→倉、動詞近→進、商品名華族→滑鼠），三種修正機制同時作用救回、"
             "真的把 50 件寫進庫存。底部拼音證據解釋原理：ASR 錯字特徵是「音同/音近但字形零重疊」，"
             "字形比對必敗、轉拼音一比即中——這就是為什麼發音容錯層有效。")
pn(s)


# ─── S13e 語音 POC · 救不回的極限（誠實交代）★ ────────────────────
s = slide_blank()
title_bar(s, "VOICE POC · 極限與對策", "救不回的也照實講：270M 的天花板")
add_text(s, MX, Inches(1.38), Inches(11.8), Inches(0.4),
         "誠實區分：同音/音近能救；但 ASR 把整詞聽成毫不相干的詞（字形+發音都差），容錯層無能為力。",
         size=13, color=GREY55)
# 左：救不回案例表
add_round(s, MX, Inches(1.95), Inches(7.15), Inches(4.15), fill=LIGHT, shadow=True)
add_text(s, MX + Inches(0.3), Inches(2.13), Inches(6.6), Inches(0.36),
         "整詞崩壞 → 救不回（FAIL）", size=14.5, bold=True, color=CORAL)
fail_rows = [
    ("瑜珈墊有貨嗎", "女藥店有貨嗎"),
    ("露營帳篷有貨嗎", "女人占房有好嗎"),
    ("北倉的滑鼠有多少", "北倉的瓦數有多少"),
    ("有缺的列出來", "有趣的列出來"),
    ("北倉進三十瓶防蚊液", "…近三十品防瘟疫"),
    ("南倉收了三十個啤酒", "…收了三十個皮"),
    ("啞鈴庫存", "10 庫存"),
]
fhy = Inches(2.62); fhh = Inches(0.475)
add_text(s, MX + Inches(0.3), fhy - Inches(0.02), Inches(3.2), Inches(0.3),
         "原句", size=11, bold=True, color=GREY77)
add_text(s, MX + Inches(3.75), fhy - Inches(0.02), Inches(3.2), Inches(0.3),
         "ASR 聽成", size=11, bold=True, color=GREY77)
for r, (o, w) in enumerate(fail_rows):
    y = fhy + Inches(0.3) + (fhh) * r
    add_text(s, MX + Inches(0.3), y, Inches(3.35), Inches(0.4),
             o, size=12, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, MX + Inches(3.75), y, Inches(3.2), Inches(0.4),
             "→ " + w, size=12, color=CORAL, anchor=MSO_ANCHOR.MIDDLE)
# 右：對策
rx = Inches(8.2)
add_round(s, rx, Inches(1.95), Inches(4.4), Inches(4.15), fill=DARK, shadow=True)
add_text(s, rx + Inches(0.32), Inches(2.13), Inches(3.8), Inches(0.36),
         "對策：人在迴路即容錯", size=14.5, bold=True, color=TEAL)
cop = [
    ("🗣️", "前端即時顯示辨識文字", "訪客看到「女藥店」明顯不對，自然會重講一次"),
    ("🎯", "不硬編同音規則硬猜", "亂編規則救個案，會誤傷別句、風險大於效益"),
    ("🔊", "音量才是主因", "小聲時摩擦音 ㄕ/ㄘ 糊掉；展場大聲對麥＝多半消失"),
    ("🔁", "展後回饋閉環", "展場問答入 journal，事後撈真實錯句補規則"),
]
for i, (ic, h, d) in enumerate(cop):
    y = Inches(2.62) + Inches(0.85) * i
    dot_icon(s, rx + Inches(0.32), y, ic, d=0.4, circle=TEAL, gcolor=DARK, gsize=12)
    add_text(s, rx + Inches(0.85), y - Inches(0.04), Inches(3.4), Inches(0.32),
             h, size=12.5, bold=True, color=WHITE)
    add_text(s, rx + Inches(0.85), y + Inches(0.26), Inches(3.45), Inches(0.5),
             d, size=11, color=GREYBB, line_spacing=1.1)
# 底部一句
add_text(s, MX, Inches(6.35), Inches(11.8), Inches(0.6),
         "分寸：容錯層只救「聽得出是同一個音」的錯——救得到的自動救、救不到的交給人重講，絕不亂猜幻覺出錯商品。",
         size=13, bold=True, color=TEALDK, line_spacing=1.2)
set_notes(s, "★極限與對策頁（誠實加分）。技術評審最怕看到只報喜不報憂——這頁專講救不回的。"
             "左表：ASR 把整詞聽成毫不相干的詞（瑜珈墊→女藥店、露營帳篷→女人占房、滑鼠→瓦數、"
             "缺→趣、防蚊液→防瘟疫），字形和發音都對不上，容錯層本來就不該硬救。右邊對策核心＝"
             "人在迴路：前端即時顯示辨識文字，訪客看到明顯錯字會自然重講——這比硬編同音規則安全得多"
             "（亂編規則救個案會誤傷別句）。另外三點：音量才是主因（正常音量下大半 FAIL 會消失）、"
             "展後靠 journal 問答記錄做回饋閉環補規則。分寸金句：救得到的自動救、救不到的交給人重講，"
             "絕不亂猜幻覺出錯商品——這正是整個系統『寧可反問、不可答錯』哲學在語音上的延伸。")
pn(s)


# ─── S14d 換輸入源測試（TTS 基準批 + 探針批）★ ────────────────────
s = slide_blank()
title_bar(s, "TESTING · 換輸入源", "測了 11 輪還有破口？因為造句的人一直是同一個")
add_text(s, MX, Inches(1.35), Inches(11.8), Inches(0.4),
         "守衛全綠、劇情批 5 輪、渲染批 5 輪——但那些句子都是同一個作者打字造的，盲點會系統性重複。",
         size=13, color=GREY55)
# 上排：兩種新輸入源
srcs = [
    ("① TTS 唸 → whisper 聽", "100 句 × 3 噪音層 = 300 次辨識",
     "產物型態跟打字完全不同：撇號縮寫、黏字、聽錯詞", TEALDK),
    ("② 刻意寫「我不會寫的句型」", "47 句探針（禮貌用語 / 口語省略 / 填充詞）",
     "換作者視角：展場訪客會客氣地問，不是打命令式", NAVY),
]
for i, (nm, sub, desc, col) in enumerate(srcs):
    x = MX + Inches(6.05) * i
    add_round(s, x, Inches(1.92), Inches(5.82), Inches(1.32), fill=LIGHT, shadow=True)
    add_text(s, x + Inches(0.28), Inches(2.04), Inches(5.3), Inches(0.34),
             nm, size=13.5, bold=True, color=col)
    add_text(s, x + Inches(0.28), Inches(2.38), Inches(5.3), Inches(0.3),
             sub, size=11, color=GREY77)
    add_text(s, x + Inches(0.28), Inches(2.7), Inches(5.3), Inches(0.44),
             desc, size=11.5, color=GREY44, line_spacing=1.2)
# 中間：TTS 三層結果
add_text(s, MX, Inches(3.42), Inches(11.8), Inches(0.32),
         "TTS 基準批結果：噪音幾乎不影響（真實賣場環境音混入）",
         size=13, bold=True, color=TEALDK)
lvl = [("沒有背景音", "92%"), ("一般人潮", "92%"), ("尖峰吵雜", "91%")]
for i, (nm, sc) in enumerate(lvl):
    x = MX + Inches(3.99) * i
    add_round(s, x, Inches(3.8), Inches(3.78), Inches(0.92), fill=TEALBG,
              line=TEAL, line_w=1.1)
    add_text(s, x + Inches(0.2), Inches(3.94), Inches(1.5), Inches(0.5), sc,
             font=FONT_EN, size=24, bold=True, color=TEALDK)
    add_text(s, x + Inches(1.75), Inches(3.98), Inches(1.9), Inches(0.42),
             nm, size=11.5, color=GREY44, anchor=MSO_ANCHOR.MIDDLE)
# 下方：抓到什麼
add_text(s, MX, Inches(4.92), Inches(11.8), Inches(0.32),
         "抓到的破口——**大多數打字訪客也會遇到**，只是先前造句時沒想到".replace("**", ""),
         size=13, bold=True, color=TEALDK)
finds = [
    ("撇號縮寫", "what's in central warehouse…",
     "ASR 聽對、LLM 判對，卻被防幻覺閘門當「陌生商品」清掉 → 全店概覽", CORAL),
    ("禮貌用語", "could you tell me the earphone stock",
     "回「查無 could 這個商品」——訪客客氣問反而失敗", CORAL),
    ("英文追問", "what about north / how about central",
     "最自然的追問講法，carry-over 詞表偏偏漏了 about", AMBER),
    ("問展示本身", "what's this demo about",
     "回熱銷榜——訪客第一句就答非所問", AMBER),
]
fy = Inches(5.3)
for i, (tag, ex, eff, col) in enumerate(finds):
    y = fy + Inches(0.44) * i
    add_rect(s, MX, y, Inches(0.06), Inches(0.38), fill=col)
    add_text(s, MX + Inches(0.22), y + Inches(0.02), Inches(1.5), Inches(0.32),
             tag, size=11.5, bold=True, color=col)
    add_text(s, MX + Inches(1.85), y + Inches(0.02), Inches(3.7), Inches(0.32),
             ex, font=FONT_EN, size=10.5, color=GREY44)
    add_text(s, MX + Inches(5.7), y + Inches(0.02), Inches(6.1), Inches(0.32),
             eff, size=11, color=GREY55)
add_round(s, MX, Inches(7.08), Inches(11.87), Inches(0.34), fill=DARK)
add_rich(s, MX + Inches(0.35), Inches(7.1), Inches(11.2), Inches(0.3),
         [[{"text": "結論  ", "size": 11.5, "bold": True, "color": TEAL},
           {"text": "不是「測不夠多輪」，是「輸入源不夠多樣」——換一個產生源，"
                    "立刻抓到 11 輪都沒碰到的類型。",
            "size": 11, "color": GREYBB}]],
         anchor=MSO_ANCHOR.MIDDLE)
set_notes(s, "★這頁回答一個很自然的質疑：「測了 11 輪、守衛全綠，怎麼還有破口？」"
             "答案是——**不是輪數不夠，是輸入源不夠多樣**。前面 11 輪的句子全是同一個作者"
             "（我）打字造的，作者的盲點會系統性重複，再跑 20 輪同樣方式也抓不到。"
             "所以換了兩個產生源：①**TTS 唸出來、whisper 聽回去**——產物型態跟打字完全不同"
             "（撇號縮寫、黏字、聽錯詞）；②**刻意寫「我自己不會寫的句型」**——換作者視角，"
             "展場訪客會客氣地問（could you tell me…）而不是打命令式（earphone stock）。"
             "TTS 基準批跑 100 句 × 3 噪音層共 300 次辨識，用的是真實賣場環境音混入："
             "乾淨 92%、一般展場 92%、尖峰吵雜 91%——**噪音幾乎不影響**，whisper 對環境音"
             "的韌性比預期好。抓到的破口裡最值得講的是撇號那個：ASR **完全聽對**、LLM 也"
             "**完全判對**（keyword=mouse、warehouse=central 都抽對了），卻被防幻覺閘門把"
             "`what's` 當成「庫裡沒有的商品修飾詞」，於是清掉正確的 keyword，回了全店概覽。"
             "根因是閘門剝標點時只剝頭尾，撇號在字中間剝不掉。**這類破口打字訪客也會遇到**"
             "——撇號、禮貌用語、what about 追問都是，只是先前造句時沒想到。真正語音專屬的"
             "只有黏字（sunheadstock）和聽錯詞（mops→mobs）那兩類，那兩句實測**沒有可用訊號**"
             "能區分正確與誤配（mobs 最像的是 mouse 不是 mop），硬修會更糟，誠實留給補訓語料。")
pn(s)


# ─── S14e 真人語音實測（誠實面對落差）★ ────────────────────────────
s = slide_blank()
title_bar(s, "VOICE · 真人實測",
          "合成音 92%，真人首測 55%——合成音會嚴重高估（下頁為補滿 100 句的現況）")
add_text(s, MX, Inches(1.35), Inches(11.8), Inches(0.4),
         "先前所有英文語音數據都來自 TTS。這次請真人（非母語者）先錄 38 句抽樣（後補滿 100 句，見下頁），"
         "用展場實際會部署的模型與麥克風測完整語音鏈。",
         size=13, color=GREY55)

# 三段對照
steps = [
    ("TTS 合成音", "92%", "US 腔 99 句 · 咬字標準、無語速變化", GREY77),
    ("真人（非母語）", "55%", "台灣腔首測 38 句抽樣（後補滿100句→下頁） · 同麥克風同模型", CORAL),
    ("＋ ASR 容錯層", "79%", "**同一批音檔**，沒重錄、只加規則", TEALDK),
]
for i, (nm, pct, desc, col) in enumerate(steps):
    x = MX + Inches(4.02) * i
    add_round(s, x, Inches(1.95), Inches(3.8), Inches(1.5), fill=LIGHT, shadow=True)
    add_text(s, x + Inches(0.25), Inches(2.08), Inches(3.3), Inches(0.32),
             nm, size=13, bold=True, color=DARK)
    add_text(s, x + Inches(0.25), Inches(2.4), Inches(3.3), Inches(0.6),
             pct, font=FONT_EN, size=34, bold=True, color=col)
    add_text(s, x + Inches(0.25), Inches(3.02), Inches(3.3), Inches(0.38),
             desc.replace("**", ""), size=10.5, color=GREY55, line_spacing=1.15)

add_text(s, MX, Inches(3.62), Inches(11.8), Inches(0.32),
         "為什麼落差這麼大——這正是「合成音會嚴重高估」的又一次驗證",
         size=13, bold=True, color=TEALDK)

gaps = [
    ("詞尾被吞掉", "shipped → shed｜send → sen｜received → receive",
     "非母語者通病，母語者也有只是頻率低", TEALDK),
    ("連音黏成一詞", "sun hat → some heat｜sun hats → some headers",
     "母語者講快時更容易發生", AMBER),
    ("整詞被替換", "trash bags → trespass life｜mop → monk's",
     "無可用訊號可修——硬修會誤傷正常查詢", CORAL),
]
gy = Inches(4.0)
for i, (tag, ex, note, col) in enumerate(gaps):
    y = gy + Inches(0.52) * i
    add_rect(s, MX, y, Inches(0.06), Inches(0.44), fill=col)
    add_text(s, MX + Inches(0.24), y + Inches(0.02), Inches(2.0), Inches(0.36),
             tag, size=12, bold=True, color=col)
    add_text(s, MX + Inches(2.4), y + Inches(0.02), Inches(4.6), Inches(0.36),
             ex, font=FONT_EN, size=10.5, color=GREY44)
    add_text(s, MX + Inches(7.3), y + Inches(0.02), Inches(4.5), Inches(0.36),
             note, size=11, color=GREY55)

add_round(s, MX, Inches(5.72), Inches(11.87), Inches(1.5), fill=DARK, shadow=True)
add_rich(s, MX + Inches(0.4), Inches(5.86), Inches(11.1), Inches(0.44),
         [[{"text": "容錯層買到什麼  ", "size": 13, "bold": True, "color": TEAL},
           {"text": "55% → 79%（+24 個百分點），完全沒有重錄", "size": 13,
            "bold": True, "color": WHITE},
           {"text": "——只在 ASR 出口加三類規則：動詞詞尾、商品名固定錯法、倉別。",
            "size": 11.5, "color": GREYBB}]],
         anchor=MSO_ANCHOR.MIDDLE)
add_rich(s, MX + Inches(0.4), Inches(6.34), Inches(11.1), Inches(0.44),
         [[{"text": "誠實揭露  ", "size": 13, "bold": True, "color": AMBER},
           {"text": "79% 是「重錄挑最好那次」的結果，一次命中率更低",
            "size": 12, "bold": True, "color": WHITE},
           {"text": "——展場只有一次機會。母語者預估 75-85%。",
            "size": 11.5, "color": GREYBB}]],
         anchor=MSO_ANCHOR.MIDDLE)
add_rich(s, MX + Inches(0.4), Inches(6.82), Inches(11.1), Inches(0.34),
         [[{"text": "設計取捨  ", "size": 12, "bold": True, "color": TEAL},
           {"text": "答不出來時誠實反問、不亂猜——錯的時候不會給出錯誤數字。",
            "size": 11.5, "color": GREYBB}]],
         anchor=MSO_ANCHOR.MIDDLE)
set_notes(s, "★真人語音實測頁（誠實面對落差，技術評審會問「你們測過真人嗎」）。"
             "**這是英文版第一份真人語音數據**——先前所有英文語音數字都來自 TTS 合成音。"
             "三段對照：①TTS 92%（US 腔 99 句）②真人非母語 55%（台灣腔 38 句，同一支 C930 "
             "麥克風、同一顆 whisper small-q5_0 模型、走完整語音鏈到倉管判定）"
             "③加 ASR 容錯層後 **79%**——關鍵是**同一批音檔、沒有重錄**，只在 ASR 出口加規則。"
             "**TTS 高估了 37 個百分點**，這與中文版的經驗一致（當時合成音 clean 100%、"
             "真人首測只有 35/52）——所以這個專案的鐵則是：合成音只能當下限篩檢，"
             "**掛了肯定不行，過了不代表可用**。三類錯法要分清楚：詞尾被吞（shipped→shed）"
             "是非母語者通病但母語者也有；連音（sun hat→some heat）母語者講快時更容易發生；"
             "整詞被替換（trash bags→trespass life）則**沒有可用訊號能修**，硬修會誤傷正常查詢，"
             "誠實留著。容錯層只收「有規律且驗證過不誤傷」的三類，並在真人音檔上驗證 +24%。"
             "⚠️ 評審若追問可靠度，要主動說明：79% 是 user 重錄挑最好那次的結果，**一次命中率"
             "更低**；展場訪客只有一次機會。母語者預估 75-85%（TTS 測五個腔調 GB93/US92/"
             "AU90/IN90/SG77，但 TTS 不等於真人）。最後一句是設計取捨：答不出來時系統會"
             "**誠實反問**而不是亂猜，所以錯的時候不會給出錯誤數字——這比硬湊一個答案安全。")
pn(s)


# ─── S14f 真人 100 句最終實測 ★（2026-08-02）────────────────────
s = slide_blank()
title_bar(s, "VOICE · 真人 100 句實測", "真人錄 100 句，混入賣場人潮音再測一次")
add_text(s, MX, Inches(1.32), Inches(11.8), Inches(0.4),
         "用真人錄的 100 句，分別在安靜、一般人潮、尖峰吵雜三種背景音下測試——"
         "同一批錄音，只是把背景音混進去，沒有重錄。",
         size=13, color=GREY55)

tiers = [
    ("安靜環境", "80%", "沒有背景音", TEALDK),
    ("一般人潮", "82%", "賣場人聲、腳步聲", TEAL),
    ("尖峰吵雜", "72%", "人潮加倍的最壞情況", GREY77),
]
for i, (nm, pct, desc, col) in enumerate(tiers):
    x = MX + Inches(4.02) * i
    add_round(s, x, Inches(1.92), Inches(3.8), Inches(1.5), fill=LIGHT, shadow=True)
    add_text(s, x + Inches(0.25), Inches(2.05), Inches(3.3), Inches(0.32),
             nm, size=13, bold=True, color=DARK)
    add_text(s, x + Inches(0.25), Inches(2.37), Inches(3.3), Inches(0.6),
             pct, font=FONT_EN, size=34, bold=True, color=col)
    add_text(s, x + Inches(0.25), Inches(2.99), Inches(3.3), Inches(0.38),
             desc, size=10.5, color=GREY55, line_spacing=1.15)

add_text(s, MX, Inches(3.58), Inches(11.8), Inches(0.32),
         "有背景音不但沒變差，還略好一點——展場的環境音在容忍範圍內",
         size=13, bold=True, color=TEALDK)

rows = [
    ("模型聽對的只有四成",
     "但那是「一個字都不能錯」的算法；實際每 10 個字才錯 1 個", DARK),
    ("容錯層把它接起來",
     "自動修正聽錯的字（藍芽→藍牙、瑜伽店→瑜珈墊），答對率拉到八成", TEALDK),
    ("聽不懂時不硬猜",
     "救不回的句子會反問「請再說一次」，而不是自信地答錯", TEALDK),
    ("這批數字怎麼來的",
     "真人錄音 100 句 × 三種背景音，全部重跑一次，不做任何人工放寬", GREY55),
]
y = Inches(4.02)
for nm, desc, col in rows:
    add_round(s, MX, y, Inches(11.8), Inches(0.52), fill=WHITE, line=GREYE6)
    add_text(s, MX + Inches(0.28), y + Inches(0.04), Inches(3.5), Inches(0.30),
             nm, size=12, bold=True, color=col)
    add_text(s, MX + Inches(3.95), y + Inches(0.04), Inches(7.6), Inches(0.44),
             desc, size=11, color=GREY55, line_spacing=1.12)
    y += Inches(0.575)

add_rect(s, MX, Inches(6.98), Inches(11.8), Inches(0.44), fill=DARK)
add_rich(s, MX + Inches(0.4), Inches(7.02), Inches(11.1), Inches(0.34),
         [[{"text": "多輪對話  ", "size": 12, "bold": True, "color": TEAL},
           {"text": "代稱追問、確認落地在同一條連線下實測 100%——展場訪客走的正是這條路。",
            "size": 11.5, "color": GREYBB}]],
         anchor=MSO_ANCHOR.MIDDLE)
set_notes(s, "★這頁講「真人講話、有背景音」的實際表現——最貼近展場的數字。\n"
             "測法：user 本人錄 100 句，把賣場人潮音混進同一批錄音，分三種強度測。"
             "**沒有重錄任何一句**，變的只有背景音，所以三個數字可以直接比較。\n"
             "① 為何有背景音反而略好（82% > 80%）：差 2 句在雜訊範圍內，"
             "重點是「加了背景音沒有變差」——代表展場環境音在容忍範圍內。\n"
             "② 為何「模型聽對只有四成」卻能答對八成：四成是「一個字都不能錯」的"
             "嚴格算法，實際上每 10 個字才錯 1 個（多是藍芽/藍牙這種同音字）。"
             "容錯層自動修掉這些，答對率就上到八成。\n"
             "③ 救不回的怎麼辦：系統會反問「請再說一次」，不會硬猜。"
             "展場上訪客對「我聽不懂」的容忍度，遠高於「自信地答錯」。\n"
             "★ 若被問母語者會不會更高：**不敢保證**——我們只測過台灣腔真人與合成音，"
             "而合成音已證實會高估。母語者從未實測，這是誠實的空白。\n"
             "★ 若被問這數字可信嗎：這批是 2026-08-06 重跑的，不做任何人工放寬"
             "（先前版本曾把幾句「測法造成的假失敗」複判成通過，這次一律不做）。")
pn(s)


# ─── S14g 收斂日：換 9 個角度測 ★（2026-08-02）──────────────────
s = slide_blank()
title_bar(s, "QUALITY · 換角度測試", "同一批錄音當輸入源，換 9 個角度找破口")
add_text(s, MX, Inches(1.32), Inches(11.8), Inches(0.4),
         "守衛全綠不代表沒問題——換一個角度就抓到一批，補進守衛庫才算收口。",
         size=13, color=GREY55)

# 左：測試角度與成果
add_text(s, MX, Inches(1.86), Inches(5.7), Inches(0.3),
         "測試角度（依抓到破口數排序）", size=12, bold=True, color=DARK)
angles = [
    ("真實 ASR 錯法 / 拼字變體", "10 項", CORAL),
    ("寫入資料正確性 · 惡意輸入", "1 項", CORAL),
    ("跨查詢介面一致性", "1 項", CORAL),
    ("UI 提示句實際照打", "1 項", CORAL),
    ("前端互動（真的用滑鼠點）", "0", TEALDK),
    ("狀態污染（訪客不照劇本走）", "0", TEALDK),
]
y = Inches(2.2)
for nm, cnt, col in angles:
    add_round(s, MX, y, Inches(5.7), Inches(0.52), fill=WHITE, line=GREYE6)
    add_text(s, MX + Inches(0.24), y + Inches(0.09), Inches(4.2), Inches(0.34),
             nm, size=11.5, color=DARK)
    add_text(s, MX + Inches(4.6), y + Inches(0.09), Inches(0.9), Inches(0.34),
             cnt, font=FONT_EN, size=13, bold=True, color=col)
    y += Inches(0.6)
add_text(s, MX, y + Inches(0.06), Inches(5.7), Inches(0.34),
         "最後兩輪零破口 → 收斂訊號", size=11.5, bold=True, color=TEALDK)

# 右：最嚴重的三個破口
rx = MX + Inches(6.1)
add_text(s, rx, Inches(1.86), Inches(5.7), Inches(0.3),
         "最嚴重的三個（都會影響資料正確性）", size=12, bold=True, color=DARK)
bugs = [
    ("小數數量被抽成錯值",
     "「進 1.5 個滑鼠」開出 +5 的確認卡，整數部分完全丟失"),
    ("一句兩商品只記第一筆",
     "第二個商品默默消失、卡片毫無提示——訪客以為兩筆都記了"),
    ("寫入後查不到那筆",
     "庫存有變但進出紀錄查不到 → 訪客以為沒成功、可能重複進貨"),
]
y = Inches(2.2)
for nm, desc in bugs:
    add_round(s, rx, y, Inches(5.7), Inches(1.16), fill=LIGHT, shadow=True)
    add_text(s, rx + Inches(0.26), y + Inches(0.12), Inches(5.2), Inches(0.32),
             nm, size=12, bold=True, color=CORAL)
    add_text(s, rx + Inches(0.26), y + Inches(0.46), Inches(5.2), Inches(0.6),
             desc, size=11, color=GREY55, line_spacing=1.15)
    y += Inches(1.26)

add_text(s, rx, y + Inches(0.02), Inches(5.7), Inches(0.56),
         "共同模式：中文版有保護、英文版判準沒英文化——"
         "移植時要優先掃「數值保護」類的機制",
         size=11, bold=True, color=TEALDK, line_spacing=1.2)

add_rect(s, MX, Inches(6.98), Inches(11.8), Inches(0.44), fill=DARK)
add_rich(s, MX + Inches(0.4), Inches(7.02), Inches(11.1), Inches(0.34),
         [[{"text": "為什麼重要  ", "size": 12, "bold": True, "color": TEAL},
           {"text": "路由錯訪客看得出來，數字錯不會——所以數值正確性的破口最該優先修。",
            "size": 11.5, "color": GREYBB}]],
         anchor=MSO_ANCHOR.MIDDLE)
set_notes(s, "★收斂日測試方法頁（2026-08-02）。核心訊息：**守衛全綠不代表沒問題**，"
             "換一個角度就抓到一批。這天用同一批 100 句真人錄音當輸入源，換了九個角度："
             "真實 ASR 錯法重放、拼字變體、寫入資料正確性、惡意邊界輸入、跨查詢介面一致性、"
             "UI 提示句實際照打、多輪長對話、並發壓測、前端真實點擊、狀態污染。"
             "前四個角度各抓到破口，**最後兩輪（前端互動、狀態污染）零破口＝收斂訊號**。"
             "最嚴重的三個都跟資料正確性有關：①小數數量『進 1.5 個』開出 +5 的卡"
             "（中文版有保護、英文版判準要求中文量詞所以失效）②一句兩商品只記第一筆"
             "（中文版有攔截、英文版判準是中文量詞與連接詞）③寫入後查不到那筆"
             "（熱更新只更新庫存、漏了進出紀錄，中英文版都中）。"
             "**共同模式是「中文版有保護、英文版判準沒英文化」**——移植時要優先掃"
             "數值保護類的機制（上限、負數、零、小數），因為**路由錯訪客看得出來、"
             "數字錯不會**。整天守衛庫維持 892/892 零回歸。")
pn(s)


# ─── S12 RPI5 實戰驗收 ────────────────────────────────────
s = slide_blank()
title_bar(s, "REAL HARDWARE", "不是實驗室數字：樹莓派上真的扛得住")
kpi_row(s, Inches(1.85), [
    ("33 hr", "連續運行"),
    ("1600+", "次推論"),
    ("44°C", "溫度穩定"),
    ("20–30 t/s", "速度零衰減"),
], num_size=26)
add_round(s, MX, Inches(3.4), Inches(5.75), Inches(3.1), fill=LIGHT, shadow=True)
add_text(s, MX + Inches(0.35), Inches(3.65), Inches(5.0), Inches(0.45),
         "雙平台驗收原則", size=16, bold=True, color=TEALDK)
add_text(s, MX + Inches(0.35), Inches(4.2), Inches(5.0), Inches(2.1),
         "本機（Windows）快速迭代，\n樹莓派（RPi5 CPU）最終驗收。\n\n"
         "單向規則：樹莓派過 = 過。\n首次上機就抓到本機測不出的\n平台精度差異句。",
         size=14, color=GREY44, line_spacing=1.35)
add_round(s, Inches(6.85), Inches(3.4), Inches(5.75), Inches(3.1), fill=DARK, shadow=True)
add_text(s, Inches(7.2), Inches(3.65), Inches(5.0), Inches(0.45),
         "展場穩定性設計", size=16, bold=True, color=TEAL)
items = [("📶", "離線手機熱點運行，資料不出場"),
         ("🔁", "Wi-Fi 掉線自癒 watchdog（10 秒自檢）"),
         ("♻️", "一鍵重置回乾淨快照，玩壞也回得來"),
         ("🔒", "寫入操作二次確認，防誤觸")]
for i, (ic, t) in enumerate(items):
    yy = Inches(4.2) + Inches(0.56) * i
    dot_icon(s, Inches(7.2), yy, ic, d=0.38, circle=TEAL, gcolor=DARK, gsize=12)
    add_text(s, Inches(7.72), yy + Inches(0.02), Inches(4.7), Inches(0.42),
             t, size=13, color=GREYBB, anchor=MSO_ANCHOR.MIDDLE)
set_notes(s, "強調這不是實驗室數據，是真的在樹莓派硬體上跑過的。33 小時連續、1600 次推論、"
             "溫度穩定不降速。右邊是展場特別做的穩定性設計：離線運行、斷線自癒、一鍵重置、"
             "寫入二次確認。")
pn(s)

print("S11-S12 done")


# ─── S12b 硬體路線圖（現在 CPU → 未來自研晶片）點綴 ───────────────
s = slide_blank()
title_bar(s, "ROADMAP · 硬體路線", "軟體已就緒，就等算力放大")
add_text(s, MX, Inches(1.42), Inches(11.8), Inches(0.4),
         "同一套軟體架構，換上更強的算力就能跑更大的模型——這正是晶片團隊的下一步。",
         size=13.5, color=GREY55)
# 左：現在（實測，實色）
add_round(s, MX, Inches(2.1), Inches(5.55), Inches(3.9), fill=LIGHT, shadow=True)
add_text(s, MX + Inches(0.35), Inches(2.35), Inches(4.9), Inches(0.4),
         "現在 · 已實測", size=16, bold=True, color=TEALDK)
add_text(s, MX + Inches(0.35), Inches(2.8), Inches(4.9), Inches(0.5),
         "RPi5 CPU（純軟體）", size=17, bold=True, color=DARK)
now_pts = [("模型", "FunctionGemma 270M"),
           ("速度", "20–30 tokens/s"),
           ("品質", "六套回歸雙平台 100%"),
           ("成本", "一台樹莓派，無 GPU / 無雲端")]
for i, (k, v) in enumerate(now_pts):
    y = Inches(3.5) + Inches(0.58) * i
    add_text(s, MX + Inches(0.35), y, Inches(1.1), Inches(0.4), k, size=13,
             bold=True, color=TEAL)
    add_text(s, MX + Inches(1.45), y, Inches(3.8), Inches(0.4), v, size=13, color=GREY44)
# 中：箭頭
add_arrow(s, Inches(6.45), Inches(3.75), Inches(0.62), Inches(0.5), fill=TEAL)
# 右：未來（roadmap，虛線淺色）
rx = Inches(7.25)
_fut = add_round(s, rx, Inches(2.1), Inches(5.35), Inches(3.9), fill=WHITE,
                 line=TEAL, line_w=1.5)
add_text(s, rx + Inches(0.35), Inches(2.35), Inches(4.7), Inches(0.4),
         "下一階段 · Roadmap 目標", size=16, bold=True, color=TEAL)
add_text(s, rx + Inches(0.35), Inches(2.8), Inches(4.7), Inches(0.5),
         "RPi5 + 自研晶片加速", size=17, bold=True, color=DARK)
fut_pts = [("模型", "跳階到 3B / 7B 更大 LLM"),
           ("能力", "從「查詢路由」→ 真正的推理對話"),
           ("架構", "軟體不動，換算力即可放大"),
           ("狀態", "晶片開發中，應用已備妥")]
for i, (k, v) in enumerate(fut_pts):
    y = Inches(3.5) + Inches(0.58) * i
    add_text(s, rx + Inches(0.35), y, Inches(1.1), Inches(0.4), k, size=13,
             bold=True, color=AMBER)
    add_text(s, rx + Inches(1.45), y, Inches(3.6), Inches(0.4), v, size=13, color=GREY44)
add_text(s, MX, Inches(6.35), Inches(11.8), Inches(0.5),
         "關鍵訊息：連純 CPU 都已跑出生產級品質——軟體、測試、落地全部就緒，就等自研晶片把天花板拉高。",
         size=13, bold=True, color=TEALDK)
set_notes(s, "★硬體路線圖頁（點綴，但對晶片團隊的參展定位很重要）。左邊實色=現在已實測："
             "270M 在 RPi5 純 CPU、20–30 t/s（模擬全開時 ~20、關閉 ~30）、六套回歸雙平台 100%、極低成本。右邊白底虛線框"
             "=下一階段 roadmap 目標（明確標成目標，不假裝已達成）：加自研晶片加速 → 跑 "
             "3B/7B 更大模型 → 能力從查詢升級到真正對話。核心訊息：軟體與應用已就緒，"
             "就等晶片把算力天花板拉高。誠實區分實測與目標，保住對評審的信任。")
pn(s)


# ─── S12d 雙機交付 · 第二台 RPI5 ★（2026-08-05 認證）──────────────
s = slide_blank()
title_bar(s, "DUAL UNIT · 交付準備", "第二台 RPI5 重建完成：與主機同級認證、可交客戶")
add_text(s, MX, Inches(1.40), Inches(11.8), Inches(0.4),
         "客戶交機版與展場主機完全同構——同一份程式碼（md5 對版）、同一套守衛認證、同一條還原手冊。",
         size=13, color=GREY55)
kpi_row(s, Inches(1.98), [
    (f"{GUARD_ZH}+{GUARD_EN}", "全量守衛 中/英 100%"),
    ("26/26 ×6", "並發串線三輪零串線"),
    ("65.9°C", "壓測峰值 · 全程零降頻"),
    ("414/414", "交機快篩驗收"),
], box_h=1.3, num_size=30)
_du_pts = [
    ("開機即就緒", "開機自動歸零回乾淨基準（uptime 閘門、離線免時鐘）→ 服務自啟 → "
                   "kiosk 雙語分頁自開 → 模擬自動起跑——插電就是展示狀態"),
    ("手機動線", "切熱點 → 掃 QR → 手機直連查詢/語音——DHCP、憑證、手機版面全鏈打通"
                 "（iPhone 實測）"),
    ("重建可複製", "重建缺口 9 類全數記錄成手冊 §15「一次還原到底」檢查表——"
                   "下一台照抄指令＋9 項驗收，不再踩雷"),
    ("遠端救援", "雙機 ZeroTier 就位（10.35.219.22 / .47）——展場斷網用手機熱點即可遠端搶修"),
]
for i, (k, v) in enumerate(_du_pts):
    y = Inches(3.58) + Inches(0.80) * i
    dot_icon(s, MX, y + Inches(0.05), "●", d=0.34, circle=NAVY, gsize=10)
    add_text(s, MX + Inches(0.55), y, Inches(1.9), Inches(0.7), k, size=13.5, bold=True,
             color=NAVY)
    add_text(s, MX + Inches(2.5), y, Inches(9.3), Inches(0.74), v, size=12.5, color=GREY44)
add_text(s, MX, Inches(6.9), Inches(11.8), Inches(0.5),
         "關鍵訊息：不是「複製一台機器」，是「複製一整套可驗收的交付流程」——手冊在版控裡，第三台起照走。",
         size=13, bold=True, color=TEALDK)
set_notes(s, "★雙機交付頁（2026-08-05 認證完成）。第二台從重建到可交客戶的完整故事：全量守衛"
             f"{GUARD_ZH}+{GUARD_EN} 與 parity 跟主機同級全綠；耐力賽三輪（並發串線 26/26×6、長對話、劇情批）"
             "零異常零降頻；散熱裝好後待機 42°C、壓測 65.9°C。重建過程實際踩出 9 類缺口"
             "（字型/輸入法/DHCP/桌面圖示鏈/crontab/瀏覽器政策等），全部固化成還原手冊 §15"
             "『一次還原到底』——缺口清單＋照抄指令＋9 項驗收清單，這份手冊跟程式碼一起進版控。"
             "對客戶的意義：交付的不是一台調好的機器，是一套可重複、可驗收的交付流程。")
pn(s)


# ═══════════════════════════════════════════════════════════
# ─── S13 總結（深底）─────────────────────────────────────
s = slide_blank()
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, fill=DARK)
add_rect(s, 0, 0, SLIDE_W, Inches(0.14), fill=TEAL)
add_text(s, MX, Inches(0.7), Inches(11.8), Inches(0.35),
         "SUMMARY", font=FONT_EN, size=13, bold=True, color=TEAL)
add_text(s, MX, Inches(1.15), Inches(11.8), Inches(0.7),
         "一句話查倉管，跑在樹莓派，品質可證明", size=28, bold=True, color=WHITE)
cards = [
    ("價值", "🎯", "自然語言問倉管，秒級回答；手機掃碼、離線可用、硬體成本極低"),
    ("技術", "🧠", "270M 小模型當路由器 + 規則層當決策者，業界邊緣 Agent 的正解"),
    ("品質", "🛡️", f"中文 {GUARD_ZH} 句雙平台 100%、英文 {GUARD_EN} 句 100%；「可靠」是數字不是祈禱"),
    ("語音", "🎙️", f"全離線 whisper；中英統一 {ASR['name']}（{ASR['lat']}/句）、同一套 runtime"),
    ("雙語", "🌐", "中英雙版同機並存，訪客點分頁切換；移植 19 類坑已歸納成方法論"),
]
y0 = Inches(1.92)
for i, (tag, ic, desc) in enumerate(cards):
    y = y0 + Inches(0.93) * i
    add_round(s, MX, y, Inches(11.87), Inches(0.79), fill=NAVY)
    add_icon_circle(s, MX + Inches(0.3), y + Inches(0.18), Inches(0.52), tag[0],
                    circle=TEAL, gcolor=DARK, gsize=16)
    add_text(s, MX + Inches(1.15), y + Inches(0.12), Inches(1.6), Inches(0.62),
             tag, size=17, bold=True, color=TEAL, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, MX + Inches(2.75), y + Inches(0.12), Inches(8.8), Inches(0.62),
             desc, size=13.5, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, MX, Inches(6.66), Inches(11.8), Inches(0.5),
         "下一步：第二台交機（全量認證已過）、展前真人語音多句實測、SEMICON 9/2–9/4 展出——開機自動歸零 / 動態模擬 / 雙語 / 手機動線全數就緒",
         size=13, color=GREYBB)
set_notes(s, "總結五張卡：價值、技術、品質、語音、雙語。收尾一句下一步。整份簡報的主軸——"
             "用最小的模型、最便宜的硬體，做到可以用數字證明的品質。"
             "品質那張卡的兩個數字要記熟：中文守衛 1122 句雙平台 100%、英文守衛 892 句 100%。")
pn(s)

# ─── 存檔 ────────────────────────────────────────────────
prs.save(OUT)
print(f"OK 已存檔 {OUT} ，共 {len(prs.slides.__iter__.__self__._sldIdLst)} 頁")
